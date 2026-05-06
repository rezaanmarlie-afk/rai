import os
import re
import csv
import io
import base64
import hashlib
from datetime import datetime, timedelta
from typing import Any, Dict, List, Optional

import requests
from dotenv import load_dotenv
from fastapi import FastAPI, Request, Query, UploadFile, File, Form
from fastapi.responses import HTMLResponse, JSONResponse, StreamingResponse
from fastapi.templating import Jinja2Templates

load_dotenv()

app = FastAPI(title="Jira Ticket Quality Compliance AI")

templates = Jinja2Templates(directory="templates")
templates.env.cache = {}

from jinja2 import Environment
from fastapi.templating import Jinja2Templates

jinja_env = Environment(
    loader=None,
    autoescape=True,
    cache_size=0,   # THIS disables caching properly
)

templates = Jinja2Templates(directory="templates", env=jinja_env)


JIRA_BASE_URL = os.getenv("JIRA_BASE_URL", "").rstrip("/")
JIRA_EMAIL = os.getenv("JIRA_EMAIL", "")
JIRA_API_TOKEN = os.getenv("JIRA_API_TOKEN", "")
DEFAULT_PROJECT = os.getenv("JIRA_PROJECT_KEY", "ARC").strip().upper()
DONE_STATUS_NAMES = {"done", "completed", "complete", "closed", "resolved", "cancelled", "canceled"}

JIRA_STORY_POINTS_FIELD = os.getenv("JIRA_STORY_POINTS_FIELD", "").strip()

SECTION_ORDER = [
    "Business Need",
    "Context",
    "Problem Statement",
    "Scope",
    "Dependencies",
    "Success Criteria",
    "Acceptance Criteria",
    "Definition of Ready",
    "Definition of Done",
    "Risks / Controls",
    "Additional Notes",
]

SECTION_PATTERNS = {
    "Business Need": [r"business need\s*:"],
    "Context": [r"context\s*:"],
    "Problem Statement": [r"problem statement\s*:"],
    "Scope": [r"scope\s*:"],
    "Dependencies": [r"dependencies\s*:"],
    "Success Criteria": [r"success criteria\s*:"],
    "Acceptance Criteria": [r"acceptance criteria\s*:"],
    "Definition of Ready": [r"definition of ready\s*:", r"\bdor\b"],
    "Definition of Done": [r"definition of done\s*:", r"\bdod\b"],
    "Risks / Controls": [r"risks?\s*/\s*controls\s*:", r"risks\s*:"],
    "Additional Notes": [r"additional notes\s*:"],
}


def jira_auth_headers() -> Dict[str, str]:
    token = f"{JIRA_EMAIL}:{JIRA_API_TOKEN}"
    encoded = base64.b64encode(token.encode()).decode()
    return {
        "Authorization": f"Basic {encoded}",
        "Accept": "application/json",
        "Content-Type": "application/json",
    }


def normalize_project_key(project_key: Optional[str]) -> str:
    key = (project_key or DEFAULT_PROJECT).strip().upper()
    return key or DEFAULT_PROJECT


def safe_json_response(response: requests.Response) -> Dict[str, Any]:
    try:
        body = response.json()
    except Exception:
        body = {"raw_text": response.text}
    return {"status_code": response.status_code, "jira_response": body}


def extract_plain_text_from_adf(adf: Any) -> str:
    if not adf:
        return ""
    texts: List[str] = []

    def walk(node: Any) -> None:
        if isinstance(node, dict):
            if node.get("type") == "text":
                texts.append(node.get("text", ""))
            for value in node.values():
                walk(value)
        elif isinstance(node, list):
            for item in node:
                walk(item)

    walk(adf)
    flattened = " ".join(part.strip() for part in texts if str(part).strip())
    return re.sub(r"\s+", " ", flattened).strip()


def build_adf_document(text: str) -> Dict[str, Any]:
    paragraphs = [p.strip() for p in (text or "").split("\n") if p.strip()] or [""]
    return {
        "type": "doc",
        "version": 1,
        "content": [
            {"type": "paragraph", "content": [{"type": "text", "text": para}]}
            for para in paragraphs
        ],
    }


def split_sentences(text: str) -> List[str]:
    if not text:
        return []
    parts = re.split(r"(?<=[.!?])\s+|\n+", text)
    return [p.strip() for p in parts if p and p.strip()]


def first_nonempty(items: List[str]) -> str:
    for item in items:
        if item and item.strip():
            return item.strip()
    return ""


def dedupe_keep_order(items: List[str], limit: int = 10) -> List[str]:
    out: List[str] = []
    seen = set()
    for item in items:
        clean = re.sub(r"\s+", " ", str(item).strip())
        if clean and clean.lower() not in seen:
            seen.add(clean.lower())
            out.append(clean)
        if len(out) >= limit:
            break
    return out


def normalize_to_list(value: Any) -> List[str]:
    if value is None:
        return []
    if isinstance(value, list):
        items = [str(x).strip() for x in value if str(x).strip()]
        return dedupe_keep_order(items, 20)
    if isinstance(value, str):
        parts = re.split(r"\n|,|;|•|- ", value)
        items = [p.strip() for p in parts if p.strip()]
        return dedupe_keep_order(items, 20)
    return [str(value).strip()]


def find_sentences_with_keywords(text: str, keywords: List[str], max_hits: int = 8) -> List[str]:
    hits: List[str] = []
    for sentence in split_sentences(text):
        lower = sentence.lower()
        if any(keyword.lower() in lower for keyword in keywords):
            hits.append(sentence)
        if len(hits) >= max_hits:
            break
    return dedupe_keep_order(hits, max_hits)


def find_section_snippet(text: str, headings: List[str], window: int = 1800) -> str:
    lower = text.lower()
    for heading in headings:
        idx = lower.find(heading.lower())
        if idx != -1:
            return text[idx: idx + window].strip()
    return ""


def extract_bullets_from_text(text: str, max_items: int = 10) -> List[str]:
    candidates = re.split(r"(?:\n|•|- |\* )", text)
    cleaned = []
    for c in candidates:
        c = re.sub(r"\s+", " ", c).strip(" :-\t\r\n")
        if len(c) > 15:
            cleaned.append(c)
    return dedupe_keep_order(cleaned, max_items)


def detect_business_context(description: str) -> Dict[str, Any]:
    section = find_section_snippet(description, ["initiative summary", "overview", "business problem", "business need", "strategic value"])
    hits = find_sentences_with_keywords(description, ["business problem", "business need", "overview", "initiative summary", "operational impact", "strategic value", "results in", "outcome"])
    if section:
        hits.insert(0, section)
    return {"found": bool(hits), "evidence": dedupe_keep_order(hits, 8)}


def detect_scope(description: str) -> Dict[str, Any]:
    section = find_section_snippet(description, ["scope", "poc scope and capabilities", "scope & use cases", "use case", "use cases"])
    hits = find_sentences_with_keywords(description, ["scope", "use case", "use cases", "capabilities", "implements", "demonstrates", "out of scope", "in scope"])
    if section:
        hits.insert(0, section)
    return {"found": bool(hits), "evidence": dedupe_keep_order(hits, 8)}


def detect_dependencies(description: str, explicit_dependencies: Any = None) -> Dict[str, Any]:
    evidence = []
    methods = []
    explicit_list = normalize_to_list(explicit_dependencies)
    if explicit_list:
        evidence.extend(explicit_list)
        methods.append("explicit_dependencies_field")
    section = find_section_snippet(description, ["inputs & dependencies", "dependencies"])
    if section:
        evidence.append(section)
        methods.append("dependency_section")
    hits = find_sentences_with_keywords(description, ["dependency", "dependencies", "depends on", "required environments", "access and credentials", "vendor dependencies", "approvals", "approval", "security", "firewall", "access"])
    if hits:
        evidence.extend(hits)
        methods.append("dependency_mentions")
    return {"found": bool(evidence), "methods": dedupe_keep_order(methods, 5), "evidence": dedupe_keep_order(evidence, 8)}


def detect_acceptance(description: str, explicit_ac: Any = None) -> Dict[str, Any]:
    evidence = []
    methods = []
    explicit_list = normalize_to_list(explicit_ac)
    if explicit_list:
        evidence.extend(explicit_list)
        methods.append("explicit_acceptance_criteria_field")
    lower = description.lower()
    if all(token in lower for token in ["given", "when", "then"]):
        evidence.extend(find_sentences_with_keywords(description, ["given", "when", "then"]))
        methods.append("given_when_then")
    for headings, method in [
        (["definition of ready", "dor"], "definition_of_ready"),
        (["definition of done", "dod"], "definition_of_done"),
        (["success criteria"], "success_criteria"),
        (["acceptance criteria"], "acceptance_criteria_heading"),
        (["ready when"], "ready_when"),
        (["done when"], "done_when"),
    ]:
        section = find_section_snippet(description, headings)
        if section:
            evidence.append(section)
            methods.append(method)
    return {"found": bool(evidence), "methods": dedupe_keep_order(methods, 8), "evidence": dedupe_keep_order(evidence, 8)}


def detect_success_criteria(description: str) -> Dict[str, Any]:
    section = find_section_snippet(description, ["success criteria", "successful if", "success criteria & measurement"])
    hits = find_sentences_with_keywords(description, ["success criteria", "successful if", "measured", "reduction", "mttr", "effort", "threshold", "go/no-go", "improved"])
    if section:
        hits.insert(0, section)
    return {"found": bool(hits), "evidence": dedupe_keep_order(hits, 8)}


def detect_risk(description: str) -> Dict[str, Any]:
    section = find_section_snippet(description, ["non-functional & risk", "quality & safety", "risks", "risk"])
    hits = find_sentences_with_keywords(description, ["risk", "risks", "mitigation", "rollback", "constraints", "security", "safe", "guardrails", "negative-path"])
    if section:
        hits.insert(0, section)
    return {"found": bool(hits), "evidence": dedupe_keep_order(hits, 8)}


def detect_dor_dod(description: str) -> Dict[str, Any]:
    dor = find_section_snippet(description, ["definition of ready", "dor"])
    dod = find_section_snippet(description, ["definition of done", "dod"])
    return {"dor_found": bool(dor), "dod_found": bool(dod), "dor_evidence": [dor] if dor else [], "dod_evidence": [dod] if dod else []}


def guess_problem_statement(summary: str, description: str) -> str:
    candidates = []
    section = find_section_snippet(description, ["business problem", "problem statement"])
    if section:
        candidates.append(section)
    candidates.extend(find_sentences_with_keywords(description, ["manual", "tribal knowledge", "risk", "mttr", "inconsistent", "dependency", "slow", "current", "results in", "problem"]))
    if candidates:
        return dedupe_keep_order(candidates, 3)[0]
    if summary:
        return f"The current delivery or operational process related to '{summary}' needs a clearer, safer, and more repeatable execution path."
    return "The current process requires clearer business context, execution detail, and measurable delivery outcomes."


def infer_scope_points(summary: str, description: str) -> List[str]:
    section = find_section_snippet(description, ["scope", "poc scope and capabilities", "scope & use cases"])
    points = []
    if section:
        points.extend(extract_bullets_from_text(section, 8))
    points.extend(find_sentences_with_keywords(description, ["use case", "use cases", "implements", "demonstrates", "workflow", "execution", "integration", "validation"]))
    if not points and summary:
        points.extend([
            f"Deliver the work required for {summary}.",
            "Define the target systems, workflow, and expected operational outcome.",
            "Clarify what is in scope and what is out of scope for this ticket.",
        ])
    return dedupe_keep_order(points, 6)


def infer_dependency_points(summary: str, description: str) -> List[str]:
    points = detect_dependencies(description)["evidence"]
    if not points:
        if "firewall" in summary.lower():
            points.extend([
                "Firewall rule approval and network security validation are required.",
                "Source, destination, port, and protocol details must be confirmed before implementation.",
            ])
        if "access" in summary.lower():
            points.append("Required system access, credentials, and approvals must be available before execution.")
    if not points:
        points.extend([
            "Required systems, approvals, credentials, and stakeholder inputs must be identified before execution.",
            "Any upstream or downstream system dependency must be documented and confirmed.",
        ])
    return dedupe_keep_order(points, 6)


def infer_success_points(summary: str, description: str) -> List[str]:
    points = detect_success_criteria(description)["evidence"]
    if not points:
        if "firewall" in summary.lower():
            points.extend([
                "The change is successful when connectivity is permitted only for the approved source/destination path.",
                "Validation evidence confirms that the target service is reachable and policy-compliant.",
            ])
        else:
            points.extend([
                "The work is successful when the agreed business outcome is achieved and evidence is recorded.",
                "Execution results must be repeatable, auditable, and validated by stakeholders.",
            ])
    return dedupe_keep_order(points, 5)


def infer_acceptance_criteria(summary: str, description: str) -> List[str]:
    acceptance = detect_acceptance(description)
    if acceptance["found"] and any("given" in e.lower() and "when" in e.lower() and "then" in e.lower() for e in acceptance["evidence"]):
        return dedupe_keep_order(acceptance["evidence"], 5)
    summary_lower = summary.lower()
    if "firewall" in summary_lower:
        return [
            "Given approved source and destination details, when the firewall rule is implemented, then the permitted application servers can reach the approved endpoints successfully.",
            "Given an invalid or incomplete access request, when implementation is attempted, then the change is stopped and the missing prerequisite is identified.",
            "Given the rule is implemented, when validation is performed, then connectivity evidence and security approval are recorded.",
        ]
    if "access" in summary_lower:
        return [
            "Given the required approvals and credentials are available, when access is configured, then the intended users or systems can access the required platform successfully.",
            "Given missing or invalid prerequisites, when the change is attempted, then the process is halted and the gap is reported.",
            "Given the change is complete, when validation is performed, then access works as expected and evidence is retained.",
        ]
    if "poc" in summary_lower or "proof of concept" in description.lower():
        return [
            "Given the agreed scope and environments are available, when the POC is executed, then the target use cases run successfully end to end.",
            "Given the solution output is produced, when results are reviewed, then measurable evidence of effort reduction, consistency, or improved turnaround is available.",
            "Given the POC completes, when stakeholders assess the outcome, then a clear go/no-go or refine recommendation is documented.",
        ]
    return [
        "Given the required prerequisites are available, when the work is executed, then the intended technical outcome is achieved successfully.",
        "Given missing or invalid prerequisites, when execution is attempted, then the issue is rejected or controlled safely.",
        "Given implementation is complete, when validation is performed, then evidence confirms the expected result and stakeholder sign-off can proceed.",
    ]


def infer_dor_points(summary: str, description: str) -> List[str]:
    dor = detect_dor_dod(description)
    points = extract_bullets_from_text(" ".join(dor["dor_evidence"]), 8) if dor["dor_found"] else []
    if not points:
        points = [
            "Scope and intended outcome are agreed.",
            "Required systems, access, approvals, and dependencies are identified.",
            "Relevant stakeholders, owners, and validation approach are confirmed.",
        ]
    return dedupe_keep_order(points, 6)


def infer_dod_points(summary: str, description: str) -> List[str]:
    dod = detect_dor_dod(description)
    points = extract_bullets_from_text(" ".join(dod["dod_evidence"]), 8) if dod["dod_found"] else []
    if not points:
        points = [
            "Implementation is completed and validated.",
            "Evidence of successful execution is captured.",
            "Known issues, constraints, and stakeholder sign-off are documented.",
        ]
    return dedupe_keep_order(points, 6)


def evaluate_ticket(issue_data: Dict[str, Any]) -> Dict[str, Any]:
    summary = (issue_data.get("summary") or "").strip()
    description = (issue_data.get("description") or "").strip()
    issue_type = (issue_data.get("issue_type") or issue_data.get("issueType") or "").strip()
    labels = issue_data.get("labels") or []
    components = issue_data.get("components") or []
    priority = (issue_data.get("priority") or "").strip()
    dependencies = issue_data.get("dependencies") or []
    acceptance_criteria = issue_data.get("acceptance_criteria") or []

    evidence = {
        "business_context": detect_business_context(description),
        "scope": detect_scope(description),
        "dependencies": detect_dependencies(description, dependencies),
        "acceptance": detect_acceptance(description, acceptance_criteria),
        "success_criteria": detect_success_criteria(description),
        "risk": detect_risk(description),
        "dor_dod": detect_dor_dod(description),
    }

    score = 100
    issues: List[str] = []
    strengths: List[str] = []
    if len(summary) < 8:
        score -= 20; issues.append("Summary is too short.")
    elif len(summary) < 18:
        score -= 8; issues.append("Summary is present but not very specific.")
    else:
        strengths.append("Summary is reasonably specific.")

    if len(description) == 0:
        score -= 35; issues.append("Description is missing.")
    elif len(description) < 40:
        score -= 20; issues.append("Description is too short.")
    elif len(description) < 120:
        score -= 8; issues.append("Description exists but could use more delivery detail.")
    else:
        strengths.append("Description provides usable detail.")

    if not issue_type:
        score -= 5; issues.append("Issue type is missing.")
    else:
        strengths.append(f"Issue type is set to {issue_type}.")

    if evidence["business_context"]["found"]: strengths.append("Business context is present.")
    else: score -= 8; issues.append("Business context is not clearly stated.")

    if evidence["scope"]["found"]: strengths.append("Scope or use case detail is present.")
    else: score -= 6; issues.append("Scope or use case detail is not clearly stated.")

    if evidence["dependencies"]["found"]: strengths.append("Dependencies are documented or mentioned.")
    else: score -= 10; issues.append("Dependencies are not documented.")

    acceptance_found = evidence["acceptance"]["found"]
    dor_found = evidence["dor_dod"]["dor_found"]
    dod_found = evidence["dor_dod"]["dod_found"]
    success_found = evidence["success_criteria"]["found"]
    if acceptance_found or dor_found or dod_found or success_found:
        strengths.append("Acceptance or readiness/completion criteria are present.")
    else:
        score -= 20; issues.append("Acceptance criteria, DoR, DoD, or success criteria are missing.")

    if labels: strengths.append("Labels are present.")
    else: score -= 5; issues.append("Labels are missing.")
    if components: strengths.append("Components are present.")
    else: score -= 5; issues.append("Components are missing.")
    if priority: strengths.append(f"Priority is set to {priority}.")
    else: score -= 3; issues.append("Priority is missing.")
    if evidence["risk"]["found"]: strengths.append("Risk, safety, or mitigation content is present.")

    maturity_bonus = (4 if dor_found else 0) + (4 if dod_found else 0) + (3 if success_found else 0) + (2 if evidence["business_context"]["found"] and evidence["scope"]["found"] else 0)
    score = max(0, min(score + maturity_bonus, 100))
    if score >= 80: status, compliance = "PASS", "Strong"
    elif score >= 60: status, compliance = "PARTIAL", "Moderate"
    else: status, compliance = "FAIL", "Weak"
    return {"score": score, "status": status, "compliance": compliance, "issues": issues, "strengths": strengths, "evidence": evidence}


def merge_human_input_into_rewrite(summary: str, existing_description: str, issue_type: str, human_input: Dict[str, Any]) -> Dict[str, Any]:
    business_objective = (human_input.get("business_objective") or "").strip()
    problem_statement = (human_input.get("problem_statement") or "").strip()
    systems_involved = normalize_to_list(human_input.get("systems_involved"))
    dependencies_input = normalize_to_list(human_input.get("dependencies"))
    risks_input = normalize_to_list(human_input.get("risks"))
    success_input = normalize_to_list(human_input.get("success_criteria"))
    validation_steps = normalize_to_list(human_input.get("validation_steps"))
    stakeholders = normalize_to_list(human_input.get("stakeholders"))
    approvals = normalize_to_list(human_input.get("approvals_needed"))
    additional_notes = (human_input.get("additional_notes") or "").strip()

    business_context = detect_business_context(existing_description)
    business_text = business_objective or first_nonempty(business_context["evidence"]) or f"This work item addresses the need for {summary}."
    problem_text = problem_statement or guess_problem_statement(summary, existing_description)

    scope_points = infer_scope_points(summary, existing_description)
    if systems_involved:
        scope_points = dedupe_keep_order(scope_points + [f"Involved system/domain: {x}" for x in systems_involved], 8)
    dependency_points = dedupe_keep_order(infer_dependency_points(summary, existing_description) + dependencies_input + approvals, 8)
    success_points = dedupe_keep_order(infer_success_points(summary, existing_description) + success_input + validation_steps, 8)
    acceptance_points = infer_acceptance_criteria(summary, existing_description)
    if validation_steps:
        acceptance_points = dedupe_keep_order(acceptance_points + [f"Given implementation is completed, when validation is performed, then {x}." for x in validation_steps[:2]], 6)
    dor_points = infer_dor_points(summary, existing_description)
    if stakeholders:
        dor_points = dedupe_keep_order(dor_points + [f"Stakeholder or owner identified: {x}" for x in stakeholders], 8)
    if approvals:
        dor_points = dedupe_keep_order(dor_points + [f"Required approval identified: {x}" for x in approvals], 8)
    dod_points = infer_dod_points(summary, existing_description)
    if validation_steps:
        dod_points = dedupe_keep_order(dod_points + [f"Validation completed: {x}" for x in validation_steps], 8)
    risk_points = dedupe_keep_order(detect_risk(existing_description)["evidence"] + risks_input, 8) or ["Execution must be controlled safely, with required approvals and validation evidence retained."]

    notes_block = f"\nAdditional Notes:\n{additional_notes}" if additional_notes else ""
    rewritten = f"""Business Need:
{summary}

Context:
{business_text}

Problem Statement:
{problem_text}

Scope:
{chr(10).join(f'- {p}' for p in scope_points[:6])}

Dependencies:
{chr(10).join(f'- {p}' for p in dependency_points[:6])}

Success Criteria:
{chr(10).join(f'- {p}' for p in success_points[:5])}

Acceptance Criteria:
{chr(10).join(f'- {p}' for p in acceptance_points[:5])}

Definition of Ready:
{chr(10).join(f'- {p}' for p in dor_points[:6])}

Definition of Done:
{chr(10).join(f'- {p}' for p in dod_points[:6])}

Risks / Controls:
{chr(10).join(f'- {p}' for p in risk_points[:5])}{notes_block}""".strip()
    return {"summary": summary, "issue_type": issue_type, "rewritten_description": rewritten, "generated_sections": {"business_need": business_text, "problem_statement": problem_text, "scope": scope_points, "dependencies": dependency_points, "success_criteria": success_points, "acceptance_criteria": acceptance_points, "definition_of_ready": dor_points, "definition_of_done": dod_points, "risks": risk_points, "additional_notes": additional_notes}}


def generate_smart_rewrite(summary: str, description: str, issue_type: str = "Story") -> Dict[str, Any]:
    return merge_human_input_into_rewrite(summary, description, issue_type, {})


def clean_text_noise(text: str) -> str:
    if not text:
        return ""
    replacements = [
        (r"\bDependencies:\s*-\s*Dependencies:\s*", "Dependencies:\n- "),
        (r"\bSuccess Criteria:\s*-\s*Success Criteria:\s*", "Success Criteria:\n- "),
        (r"\bAcceptance Criteria:\s*-\s*Acceptance Criteria:\s*", "Acceptance Criteria:\n- "),
        (r"\bDefinition of Ready:\s*-\s*Definition of Ready\b", "Definition of Ready:"),
        (r"\bDefinition of Done:\s*-\s*Definition of Done\b", "Definition of Done:"),
        (r"\bRisks?\s*/\s*Controls:\s*-\s*Risks?\s*/\s*Controls:\s*", "Risks / Controls:\n- "),
        (r"-\s*-\s*", "- "),
        (r"\s+\n", "\n"),
        (r"\n{3,}", "\n\n"),
        (r"[ \t]{2,}", " "),
    ]
    out = text
    for pattern, repl in replacements:
        out = re.sub(pattern, repl, out, flags=re.IGNORECASE)
    return out.strip()


def split_into_sections(raw_text: str) -> Dict[str, List[str]]:
    text = clean_text_noise(raw_text)
    sections: Dict[str, List[str]] = {name: [] for name in SECTION_ORDER}
    heading_matches = []
    for section_name, patterns in SECTION_PATTERNS.items():
        for pattern in patterns:
            for m in re.finditer(pattern, text, flags=re.IGNORECASE):
                heading_matches.append((m.start(), m.end(), section_name))
    heading_matches.sort(key=lambda x: x[0])
    if not heading_matches:
        sections["Context"].append(text)
        return sections
    for idx, (start, end, section_name) in enumerate(heading_matches):
        next_start = heading_matches[idx + 1][0] if idx + 1 < len(heading_matches) else len(text)
        chunk = text[end:next_start].strip(" \n:-")
        if chunk:
            sections[section_name].append(chunk)
    return sections


def normalize_section_items(section_name: str, chunks: List[str]) -> List[str]:
    items: List[str] = []
    for chunk in chunks:
        if not chunk:
            continue
        if section_name in {"Scope", "Dependencies", "Success Criteria", "Acceptance Criteria", "Definition of Ready", "Definition of Done", "Risks / Controls"}:
            extracted = extract_bullets_from_text(chunk, 20)
            items.extend(extracted if extracted else split_sentences(chunk))
        else:
            items.extend(split_sentences(chunk) if section_name != "Context" else [chunk])
    cleaned_items = []
    for item in items:
        item = re.sub(r"^\s*[-•*]\s*", "", str(item)).strip()
        item = re.sub(r"\s+", " ", item).strip()
        if len(item) >= 6:
            cleaned_items.append(item)
    return dedupe_keep_order(cleaned_items, 12)


def choose_best_summary(summary: str, description: str) -> str:
    summary = (summary or "").strip()
    if len(summary) >= 20 and any(word in summary.lower() for word in ["platform", "automation", "evaluation", "engine", "system"]):
        return summary
    if "rfp" in summary.lower() or "evaluator" in summary.lower():
        return "AI-driven RFP Evaluation Platform to Automate Vendor Scoring and Reduce Technical Review Effort"
    if summary:
        return summary
    derived = first_nonempty(split_sentences(description))
    return derived[:120] if derived else "Normalized Jira Ticket"


def auto_clean_and_normalize_ticket(summary: str, description: str, issue_type: str = "Story") -> Dict[str, Any]:
    sections = split_into_sections(description or "")
    business_need = normalize_section_items("Business Need", sections["Business Need"])
    context = normalize_section_items("Context", sections["Context"])
    problem = normalize_section_items("Problem Statement", sections["Problem Statement"])
    scope = normalize_section_items("Scope", sections["Scope"])
    dependencies = normalize_section_items("Dependencies", sections["Dependencies"])
    success = normalize_section_items("Success Criteria", sections["Success Criteria"])
    acceptance = normalize_section_items("Acceptance Criteria", sections["Acceptance Criteria"])
    dor = normalize_section_items("Definition of Ready", sections["Definition of Ready"])
    dod = normalize_section_items("Definition of Done", sections["Definition of Done"])
    risks = normalize_section_items("Risks / Controls", sections["Risks / Controls"])
    notes = normalize_section_items("Additional Notes", sections["Additional Notes"])

    best_summary = choose_best_summary(summary, description)
    if not business_need: business_need = [best_summary]
    if not context:
        detected_context = detect_business_context(description)["evidence"]
        context = detected_context[:2] if detected_context else [f"This work item addresses the need for {best_summary}."]
    if not problem: problem = [guess_problem_statement(best_summary, description)]
    if not scope: scope = infer_scope_points(best_summary, description)
    if not dependencies: dependencies = infer_dependency_points(best_summary, description)
    if not success: success = infer_success_points(best_summary, description)
    if not acceptance: acceptance = infer_acceptance_criteria(best_summary, description)
    if not dor: dor = infer_dor_points(best_summary, description)
    if not dod: dod = infer_dod_points(best_summary, description)
    if not risks: risks = detect_risk(description)["evidence"] or ["Execution must be controlled safely, with required approvals and validation evidence retained."]

    normalized_description = f"""Business Need:
{business_need[0]}

Context:
{' '.join(context[:2])}

Problem Statement:
{problem[0]}

Scope:
{chr(10).join(f'- {x}' for x in dedupe_keep_order(scope, 6))}

Dependencies:
{chr(10).join(f'- {x}' for x in dedupe_keep_order(dependencies, 6))}

Success Criteria:
{chr(10).join(f'- {x}' for x in dedupe_keep_order(success, 5))}

Acceptance Criteria:
{chr(10).join(f'- {x}' for x in dedupe_keep_order(acceptance, 5))}

Definition of Ready:
{chr(10).join(f'- {x}' for x in dedupe_keep_order(dor, 6))}

Definition of Done:
{chr(10).join(f'- {x}' for x in dedupe_keep_order(dod, 6))}

Risks / Controls:
{chr(10).join(f'- {x}' for x in dedupe_keep_order(risks, 5))}""".strip()
    if notes:
        normalized_description += "\n\nAdditional Notes:\n" + "\n".join(f"- {x}" for x in dedupe_keep_order(notes, 5))
    return {"summary": best_summary, "issue_type": issue_type, "normalized_description": normalized_description, "normalized_sections": {"business_need": business_need, "context": context, "problem_statement": problem, "scope": scope, "dependencies": dependencies, "success_criteria": success, "acceptance_criteria": acceptance, "definition_of_ready": dor, "definition_of_done": dod, "risks_controls": risks, "additional_notes": notes}}


def normalize_live_issue(issue: Dict[str, Any]) -> Dict[str, Any]:
    fields = issue.get("fields", {}) or {}
    return {
        "issue_key": issue.get("key"),
        "summary": fields.get("summary") or "",
        "description": extract_plain_text_from_adf(fields.get("description")),
        "issue_type": (fields.get("issuetype") or {}).get("name", ""),
        "labels": fields.get("labels") or [],
        "components": [c.get("name", "") for c in (fields.get("components") or []) if isinstance(c, dict)],
        "priority": (fields.get("priority") or {}).get("name", ""),
        "status": (fields.get("status") or {}).get("name", ""),
        "assignee": (fields.get("assignee") or {}).get("displayName", ""),
        "reporter": (fields.get("reporter") or {}).get("displayName", ""),
        "created": fields.get("created", ""),
        "updated": fields.get("updated", ""),
        "dependencies": [],
        "acceptance_criteria": [],
    }


def summarize_bulk_results(project_key: str, jql: str, assessed_tickets: List[Dict[str, Any]]) -> Dict[str, Any]:
    total = len(assessed_tickets)
    if total == 0:
        return {"projectKey": project_key, "jql": jql, "summary": {"total_tickets": 0, "average_score": 0, "pass_count": 0, "partial_count": 0, "fail_count": 0, "missing_description_count": 0, "missing_labels_count": 0, "missing_components_count": 0, "missing_acceptance_count": 0}, "tickets": []}
    avg_score = round(sum(t["assessment"]["score"] for t in assessed_tickets) / total, 1)
    pass_count = sum(1 for t in assessed_tickets if t["assessment"]["status"] == "PASS")
    partial_count = sum(1 for t in assessed_tickets if t["assessment"]["status"] == "PARTIAL")
    fail_count = sum(1 for t in assessed_tickets if t["assessment"]["status"] == "FAIL")
    missing_description_count = sum(1 for t in assessed_tickets if not (t["normalized_issue"]["description"] or "").strip())
    missing_labels_count = sum(1 for t in assessed_tickets if len(t["normalized_issue"]["labels"]) == 0)
    missing_components_count = sum(1 for t in assessed_tickets if len(t["normalized_issue"]["components"]) == 0)
    missing_acceptance_count = sum(1 for t in assessed_tickets if not (t["assessment"]["evidence"]["acceptance"]["found"] or t["assessment"]["evidence"]["dor_dod"]["dor_found"] or t["assessment"]["evidence"]["dor_dod"]["dod_found"] or t["assessment"]["evidence"]["success_criteria"]["found"]))
    sorted_tickets = sorted(assessed_tickets, key=lambda x: (x["assessment"]["score"], x["normalized_issue"]["issue_key"]))
    return {"projectKey": project_key, "jql": jql, "summary": {"total_tickets": total, "average_score": avg_score, "pass_count": pass_count, "partial_count": partial_count, "fail_count": fail_count, "missing_description_count": missing_description_count, "missing_labels_count": missing_labels_count, "missing_components_count": missing_components_count, "missing_acceptance_count": missing_acceptance_count}, "tickets": sorted_tickets}


def is_done_or_completed_status(status_name: str) -> bool:
    return str(status_name or "").strip().lower() in DONE_STATUS_NAMES


def build_bulk_dashboard(project_key: str, jql: str, max_results: int = 0, exclude_done: bool = True) -> Dict[str, Any]:
    """Build dashboard across all matching Jira issues.

    max_results <= 0 means fetch all pages from Jira using nextPageToken.
    Jira Cloud still pages internally, so we request in chunks of 100.
    """
    fields = ["summary", "description", "status", "assignee", "issuetype", "labels", "priority", "components", "created", "updated", "reporter"]
    url = f"{JIRA_BASE_URL}/rest/api/3/search/jql"
    page_size = 100
    remaining = None if not max_results or int(max_results) <= 0 else int(max_results)
    next_page_token = None
    assessed_tickets: List[Dict[str, Any]] = []
    excluded_count = 0
    fetched_count = 0
    page_count = 0
    last_response: Dict[str, Any] = {}

    while True:
        request_size = page_size if remaining is None else max(1, min(page_size, remaining))
        params = {
            "jql": jql,
            "maxResults": request_size,
            "fields": ",".join(fields),
            "fieldsByKeys": "false",
        }
        if next_page_token:
            params["nextPageToken"] = next_page_token

        response = requests.get(url, headers=jira_auth_headers(), params=params, timeout=120)
        if response.status_code != 200:
            raise RuntimeError(str(safe_json_response(response)))

        body = response.json()
        last_response = body
        page_count += 1
        issues = body.get("issues", []) or []
        fetched_count += len(issues)

        for issue in issues:
            normalized = normalize_live_issue(issue)
            if exclude_done and is_done_or_completed_status(normalized.get("status")):
                excluded_count += 1
                continue
            assessed_tickets.append({
                "issueKey": normalized["issue_key"],
                "normalized_issue": normalized,
                "assessment": evaluate_ticket(normalized),
            })

        if remaining is not None:
            remaining -= len(issues)
            if remaining <= 0:
                break

        next_page_token = body.get("nextPageToken")
        is_last = bool(body.get("isLast", False))
        if is_last or not next_page_token or not issues:
            break

    result = summarize_bulk_results(project_key, jql, assessed_tickets)
    result["nextPageToken"] = last_response.get("nextPageToken")
    result["isLast"] = last_response.get("isLast", True)
    result["excludeDoneCompleted"] = exclude_done
    result["excluded_done_completed_count"] = excluded_count
    result["jira_fetched_count"] = fetched_count
    result["jira_pages_read"] = page_count
    result["maxResultsRequested"] = max_results
    result["allMatchingTicketsIncluded"] = (not max_results or int(max_results) <= 0)
    return result


def build_manager_dashboard_from_bulk(bulk_result: Dict[str, Any]) -> Dict[str, Any]:
    tickets = bulk_result.get("tickets", []) or []
    total = len(tickets)
    executive = {
        "total_tickets": total,
        "overall_quality_score": 0,
        "pass_count": 0,
        "partial_count": 0,
        "fail_count": 0,
        "pass_percent": 0,
        "partial_percent": 0,
        "fail_percent": 0,
        "at_risk_count": 0,
    }
    board_map: Dict[str, Dict[str, Any]] = {}
    assignee_map: Dict[str, Dict[str, Any]] = {}
    issue_breakdown: Dict[str, int] = {}
    risk_tickets: List[Dict[str, Any]] = []

    def bucket_for(container: Dict[str, Dict[str, Any]], key: str) -> Dict[str, Any]:
        if key not in container:
            container[key] = {
                "name": key,
                "total": 0,
                "score_total": 0,
                "avg_score": 0,
                "pass": 0,
                "partial": 0,
                "fail": 0,
                "missing_labels": 0,
                "missing_components": 0,
                "missing_acceptance": 0,
                "tickets": [],
            }
        return container[key]

    for t in tickets:
        ni = t.get("normalized_issue", {}) or {}
        assessment = t.get("assessment", {}) or {}
        score = int(assessment.get("score") or 0)
        status = str(assessment.get("status") or "UNKNOWN")
        issue_key = ni.get("issue_key") or t.get("issueKey") or ""
        board = issue_key.split("-")[0] if "-" in issue_key else (bulk_result.get("projectKey") or "UNKNOWN")
        assignee = ni.get("assignee") or "Unassigned"
        issues = assessment.get("issues") or []
        evidence = assessment.get("evidence") or {}
        has_acceptance = bool(
            ((evidence.get("acceptance") or {}).get("found"))
            or ((evidence.get("dor_dod") or {}).get("dor_found"))
            or ((evidence.get("dor_dod") or {}).get("dod_found"))
            or ((evidence.get("success_criteria") or {}).get("found"))
        )
        executive["pass_count"] += 1 if status == "PASS" else 0
        executive["partial_count"] += 1 if status == "PARTIAL" else 0
        executive["fail_count"] += 1 if status == "FAIL" else 0
        if score < 60:
            executive["at_risk_count"] += 1
        ticket_row = {
            "issueKey": issue_key,
            "summary": ni.get("summary") or "",
            "assignee": assignee,
            "board": board,
            "score": score,
            "status": status,
            "priority": ni.get("priority") or "",
            "main_gaps": issues[:4],
        }
        for container_key, container in [(board, board_map), (assignee, assignee_map)]:
            b = bucket_for(container, container_key)
            b["total"] += 1
            b["score_total"] += score
            b["pass"] += 1 if status == "PASS" else 0
            b["partial"] += 1 if status == "PARTIAL" else 0
            b["fail"] += 1 if status == "FAIL" else 0
            b["missing_labels"] += 1 if len(ni.get("labels") or []) == 0 else 0
            b["missing_components"] += 1 if len(ni.get("components") or []) == 0 else 0
            b["missing_acceptance"] += 1 if not has_acceptance else 0
            b["tickets"].append(ticket_row)
        for issue in issues:
            issue_breakdown[issue] = issue_breakdown.get(issue, 0) + 1
        if score < 60 or status == "FAIL":
            risk_tickets.append(ticket_row)

    if total:
        executive["overall_quality_score"] = round(sum((t.get("assessment", {}) or {}).get("score", 0) for t in tickets) / total, 1)
        executive["pass_percent"] = round(executive["pass_count"] * 100 / total, 1)
        executive["partial_percent"] = round(executive["partial_count"] * 100 / total, 1)
        executive["fail_percent"] = round(executive["fail_count"] * 100 / total, 1)
    for container in [board_map, assignee_map]:
        for v in container.values():
            v["avg_score"] = round(v["score_total"] / max(1, v["total"]), 1)
            v.pop("score_total", None)
            v["tickets"] = sorted(v["tickets"], key=lambda x: (x["score"], x["issueKey"]))
    top_performers = sorted(assignee_map.values(), key=lambda x: (-x["avg_score"], -x["total"]))[:5]
    needs_attention = sorted(assignee_map.values(), key=lambda x: (x["avg_score"], -x["total"]))[:5]
    return {
        "projectKey": bulk_result.get("projectKey"),
        "jql": bulk_result.get("jql"),
        "executive_summary": executive,
        "board_summary": sorted(board_map.values(), key=lambda x: x["name"]),
        "assignee_summary": sorted(assignee_map.values(), key=lambda x: x["avg_score"]),
        "top_performers": top_performers,
        "needs_attention": needs_attention,
        "issue_breakdown": sorted([{"issue": k, "count": v} for k, v in issue_breakdown.items()], key=lambda x: -x["count"]),
        "risk_tickets": sorted(risk_tickets, key=lambda x: x["score"]),
        "tickets": tickets,
    }


def build_manager_dashboard(project_key: str, jql: str, max_results: int, exclude_done: bool = True) -> Dict[str, Any]:
    bulk = build_bulk_dashboard(project_key, jql, max_results, exclude_done=exclude_done)
    return build_manager_dashboard_from_bulk(bulk)
def bulk_results_to_csv_rows(result: Dict[str, Any], weak_only: bool = False) -> List[List[str]]:
    rows = [["issue_key", "summary", "score", "status", "compliance", "priority", "has_description", "labels_count", "components_count", "has_acceptance_evidence", "top_issues", "top_strengths"]]
    for t in result.get("tickets", []):
        assessment = t["assessment"]
        if weak_only and assessment["status"] == "PASS":
            continue
        ni = t["normalized_issue"]
        has_acceptance = assessment["evidence"]["acceptance"]["found"] or assessment["evidence"]["dor_dod"]["dor_found"] or assessment["evidence"]["dor_dod"]["dod_found"] or assessment["evidence"]["success_criteria"]["found"]
        rows.append([t["issueKey"], ni["summary"], str(assessment["score"]), assessment["status"], assessment["compliance"], ni["priority"], "Yes" if (ni["description"] or "").strip() else "No", str(len(ni["labels"] or [])), str(len(ni["components"] or [])), "Yes" if has_acceptance else "No", " | ".join(assessment["issues"][:3]), " | ".join(assessment["strengths"][:3])])
    return rows


def csv_stream(rows: List[List[str]]) -> io.StringIO:
    output = io.StringIO()
    writer = csv.writer(output)
    writer.writerows(rows)
    output.seek(0)
    return output


def parent_field_payload(parent_issue_key: str, parent_issue_type: str, child_issue_type: str) -> Dict[str, Any]:
    parent_type = (parent_issue_type or "").strip().lower()
    child_type = (child_issue_type or "").strip().lower()
    if child_type in {"sub-task", "subtask", "sub task"}:
        return {"parent": {"key": parent_issue_key}}
    return {"parent": {"key": parent_issue_key}}


def suggest_story_points(summary: str, description: str) -> int:
    text = f"{summary} {description}".lower()
    complexity = 0
    complexity += 2 if any(k in text for k in ["integration", "api", "workflow", "engine"]) else 0
    complexity += 2 if any(k in text for k in ["multiple", "several", "across", "vendor", "comparison"]) else 0
    complexity += 2 if any(k in text for k in ["security", "approval", "data", "compliance"]) else 0
    complexity += 1 if any(k in text for k in ["dashboard", "ui", "export", "report"]) else 0
    if complexity <= 1: return 2
    if complexity <= 3: return 3
    if complexity <= 5: return 5
    if complexity <= 7: return 8
    return 13


def hierarchy_child_type(parent_issue_type: str) -> str:
    t = (parent_issue_type or "").strip().lower()
    if t == "initiative":
        return "Epic"
    if t in {"epic", "feature"}:
        return "Story"
    return "Sub-task"


def suggest_breakdown(parent_summary: str, parent_description: str, parent_issue_type: str, parent_issue_key: str = "") -> Dict[str, Any]:
    child_type = hierarchy_child_type(parent_issue_type)
    title = (parent_summary or "").strip() or "Parent work item"
    if child_type == "Epic":
        seeds = [
            "Define scope and business case",
            "Design solution architecture and workflow",
            "Build core platform capabilities",
            "Build reporting, controls, and exports",
            "Validate solution and prepare rollout",
        ]
    elif child_type == "Story":
        seeds = [
            "Define functional requirements and scoring logic",
            "Build ingestion and parsing capability",
            "Implement evaluation engine and reasoning",
            "Build dashboard, evidence, and exports",
            "Validate outputs and operational readiness",
        ]
    else:
        seeds = [
            "Design implementation approach",
            "Develop component",
            "Test and validate result",
            "Document and hand over",
        ]

    children = []
    for seed in seeds:
        summary = f"{seed} for {title}"
        description = f"""Business Need:
{summary}

Context:
This child ticket is derived from {parent_issue_type} {parent_issue_key or title}.

Problem Statement:
This work item is required to progress the parent item in a controlled and auditable way.

Scope:
- Complete the specific deliverable described in the summary
- Produce evidence of completion
- Support the parent outcome

Acceptance Criteria:
- Given prerequisites are available, when the work is executed, then the intended outcome is achieved
- Given execution is complete, when validation is performed, then evidence is recorded

Definition of Done:
- Implementation completed
- Validation completed
- Evidence captured
""".strip()
        item: Dict[str, Any] = {
            "issueType": child_type,
            "summary": summary,
            "description": description,
            "labels": [],
            "components": [],
            "priority": "Medium",
        }
        if child_type == "Story":
            item["storyPoints"] = suggest_story_points(summary, description)
        children.append(item)
    return {"parentIssueKey": parent_issue_key, "parentIssueType": parent_issue_type, "suggestedChildType": child_type, "suggestedChildren": children}


# -----------------------------
# Duplicate-safe batch creation helpers
# -----------------------------

def normalize_duplicate_text(value: str) -> str:
    value = str(value or '').lower().strip()
    value = re.sub(r'[^a-z0-9]+', ' ', value)
    return re.sub(r'\s+', ' ', value).strip()

def jira_label_safe(value: str, max_len: int = 45) -> str:
    value = str(value or '').lower().strip()
    value = re.sub(r'[^a-z0-9_-]+', '-', value)
    value = re.sub(r'-+', '-', value).strip('-')
    return value[:max_len] or 'unknown'

def breakdown_fingerprint(parent_issue_key: str, issue_type: str, summary: str) -> str:
    raw = f"{normalize_duplicate_text(parent_issue_key)}|{normalize_duplicate_text(issue_type)}|{normalize_duplicate_text(summary)}"
    return hashlib.sha1(raw.encode('utf-8')).hexdigest()[:12]

def breakdown_labels(parent_issue_key: str, issue_type: str, summary: str) -> List[str]:
    fp = breakdown_fingerprint(parent_issue_key, issue_type, summary)
    return dedupe_keep_order(['ai-breakdown', f'parent-{jira_label_safe(parent_issue_key)}', f'fp-{fp}'], 10)

def quote_jql_value(value: str) -> str:
    value = str(value or '').replace('\\', '\\\\').replace('"', '\\"')
    return f'"{value}"'

def extract_issue_parent_key(issue: Dict[str, Any]) -> str:
    fields = issue.get('fields') or {}
    parent = fields.get('parent') or {}
    return str(parent.get('key') or '').upper() if isinstance(parent, dict) else ''

def issue_matches_parent(issue: Dict[str, Any], parent_issue_key: str) -> bool:
    if not parent_issue_key:
        return True
    return extract_issue_parent_key(issue) == str(parent_issue_key or '').upper()

def search_jira_issues_for_dedup(jql: str, max_results: int = 100) -> Dict[str, Any]:
    url = f"{JIRA_BASE_URL}/rest/api/3/search/jql"
    params = {'jql': jql, 'maxResults': max_results, 'fields': 'summary,issuetype,parent,labels,status', 'fieldsByKeys': 'false'}
    response = requests.get(url, headers=jira_auth_headers(), params=params, timeout=60)
    return safe_json_response(response)

def find_existing_breakdown_ticket(project_key: str, parent_issue_key: str, issue_type: str, summary: str) -> Dict[str, Any]:
    project_key = normalize_project_key(project_key)
    issue_type = (issue_type or 'Task').strip()
    parent_issue_key = (parent_issue_key or '').strip().upper()
    target_summary_norm = normalize_duplicate_text(summary)
    fp_label = f"fp-{breakdown_fingerprint(parent_issue_key, issue_type, summary)}"
    attempts: List[Dict[str, Any]] = []

    jql_fp = f'project = {project_key} AND labels = {quote_jql_value(fp_label)} ORDER BY created DESC'
    res_fp = search_jira_issues_for_dedup(jql_fp, 50)
    attempts.append({'method': 'fingerprint_label', 'jql': jql_fp, 'status_code': res_fp.get('status_code')})
    if res_fp.get('status_code') == 200:
        for issue in (res_fp.get('jira_response') or {}).get('issues', []):
            if issue_matches_parent(issue, parent_issue_key):
                return {'found': True, 'method': 'fingerprint_label', 'existingIssueKey': issue.get('key'), 'existingSummary': ((issue.get('fields') or {}).get('summary') or ''), 'attempts': attempts}

    jql_parent = f'project = {project_key} AND issuetype = {quote_jql_value(issue_type)}'
    if parent_issue_key:
        jql_parent += f' AND parent = {quote_jql_value(parent_issue_key)}'
    jql_parent += ' ORDER BY created DESC'
    res_parent = search_jira_issues_for_dedup(jql_parent, 100)
    attempts.append({'method': 'parent_exact_summary', 'jql': jql_parent, 'status_code': res_parent.get('status_code')})
    if res_parent.get('status_code') == 200:
        for issue in (res_parent.get('jira_response') or {}).get('issues', []):
            fields = issue.get('fields') or {}
            if normalize_duplicate_text(fields.get('summary') or '') == target_summary_norm:
                return {'found': True, 'method': 'parent_exact_summary', 'existingIssueKey': issue.get('key'), 'existingSummary': fields.get('summary') or '', 'attempts': attempts}

    jql_type = f'project = {project_key} AND issuetype = {quote_jql_value(issue_type)} ORDER BY created DESC'
    res_type = search_jira_issues_for_dedup(jql_type, 100)
    attempts.append({'method': 'type_exact_summary_fallback', 'jql': jql_type, 'status_code': res_type.get('status_code')})
    if res_type.get('status_code') == 200:
        for issue in (res_type.get('jira_response') or {}).get('issues', []):
            fields = issue.get('fields') or {}
            if normalize_duplicate_text(fields.get('summary') or '') == target_summary_norm and issue_matches_parent(issue, parent_issue_key):
                return {'found': True, 'method': 'type_exact_summary_fallback', 'existingIssueKey': issue.get('key'), 'existingSummary': fields.get('summary') or '', 'attempts': attempts}
    return {'found': False, 'attempts': attempts, 'fingerprintLabel': fp_label}

def update_story_points_after_create(issue_key: str, story_points: Any, preferred_field: str = '') -> Dict[str, Any]:
    if story_points is None or story_points == '':
        return {'requested': story_points, 'updated': False, 'reason': 'No story points value supplied'}
    try:
        numeric_points = float(story_points)
    except Exception:
        return {'requested': story_points, 'updated': False, 'reason': 'Story points value is not numeric'}
    issue_key = (issue_key or '').strip().upper()
    candidate_fields: List[str] = []
    if preferred_field:
        candidate_fields.append(preferred_field)
    if JIRA_STORY_POINTS_FIELD:
        candidate_fields.append(JIRA_STORY_POINTS_FIELD)
    editmeta_response = requests.get(f"{JIRA_BASE_URL}/rest/api/3/issue/{issue_key}/editmeta", headers=jira_auth_headers(), timeout=30)
    editable_candidates = []
    if editmeta_response.status_code == 200:
        try:
            fields = (editmeta_response.json() or {}).get('fields') or {}
        except Exception:
            fields = {}
        for fid, meta in fields.items():
            name = str(meta.get('name') or '').lower()
            schema = meta.get('schema') or {}
            operations = meta.get('operations') or []
            is_number = schema.get('type') == 'number' or 'float' in str(schema.get('custom') or '').lower()
            is_story_points = name in {'story points', 'story point estimate', 'story points estimate'} or ('story' in name and 'point' in name)
            if fid.startswith('customfield_') and is_number and is_story_points and 'set' in operations:
                editable_candidates.append({'id': fid, 'name': meta.get('name'), 'schema': schema, 'operations': operations})
                candidate_fields.append(fid)
    candidate_fields = dedupe_keep_order(candidate_fields, 20)
    attempts = []
    for fid in candidate_fields:
        response = requests.put(f"{JIRA_BASE_URL}/rest/api/3/issue/{issue_key}", headers=jira_auth_headers(), json={'fields': {fid: numeric_points}}, timeout=30)
        attempt = {'field': fid, 'status_code': response.status_code}
        if response.text.strip():
            try:
                attempt['response'] = response.json()
            except Exception:
                attempt['response'] = {'raw_text': response.text}
        attempts.append(attempt)
        if response.status_code in (200, 204):
            return {'requested': numeric_points, 'updated': True, 'fieldUsed': fid, 'status_code': response.status_code, 'editableCandidates': editable_candidates, 'attempts': attempts}
    return {'requested': numeric_points, 'updated': False, 'reason': 'No Story Points candidate field accepted the update', 'editableCandidates': editable_candidates, 'attempts': attempts, 'editmetaStatus': editmeta_response.status_code}



# -----------------------------------------------------------------------------
# Document-aware ticket enhancement
# -----------------------------------------------------------------------------
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY", "").strip()
OPENAI_MODEL = os.getenv("OPENAI_MODEL", "gpt-4o-mini").strip()
SUPPORTED_DOC_EXTENSIONS = {".pdf", ".docx", ".xlsx", ".xlsm", ".pptx", ".txt", ".csv"}

def get_file_extension(filename: str) -> str:
    return os.path.splitext(filename or "")[1].lower().strip()

def extract_text_from_pdf_bytes(content: bytes) -> str:
    try:
        import pdfplumber
        parts = []
        with pdfplumber.open(io.BytesIO(content)) as pdf:
            for page_no, page in enumerate(pdf.pages, start=1):
                text = page.extract_text() or ""
                if text.strip():
                    parts.append(f"\n[PDF page {page_no}]\n{text}")
        return "\n".join(parts)
    except Exception as e:
        return f"[PDF extraction failed: {e}]"

def extract_text_from_docx_bytes(content: bytes) -> str:
    try:
        from docx import Document
        doc = Document(io.BytesIO(content))
        parts = [p.text.strip() for p in doc.paragraphs if p.text and p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                vals = [cell.text.strip() for cell in row.cells if cell.text and cell.text.strip()]
                if vals:
                    parts.append(" | ".join(vals))
        return "\n".join(parts)
    except Exception as e:
        return f"[DOCX extraction failed: {e}]"

def extract_text_from_excel_bytes(content: bytes) -> str:
    try:
        from openpyxl import load_workbook
        wb = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
        parts = []
        for ws in wb.worksheets:
            parts.append(f"\n[Excel sheet: {ws.title}]")
            for row in ws.iter_rows(values_only=True):
                vals = [str(v).strip() for v in row if v is not None and str(v).strip()]
                if vals:
                    parts.append(" | ".join(vals))
        return "\n".join(parts)
    except Exception as e:
        return f"[Excel extraction failed: {e}]"

def extract_text_from_pptx_bytes(content: bytes) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(content))
        parts = []
        for idx, slide in enumerate(prs.slides, start=1):
            slide_parts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text and shape.text.strip():
                    slide_parts.append(shape.text.strip())
            if slide_parts:
                parts.append(f"\n[PowerPoint slide {idx}]\n" + "\n".join(slide_parts))
        return "\n".join(parts)
    except Exception as e:
        return f"[PPTX extraction failed: {e}]"

def extract_text_from_uploaded_file(filename: str, content: bytes) -> Dict[str, Any]:
    ext = get_file_extension(filename)
    if ext == ".pdf":
        text = extract_text_from_pdf_bytes(content)
    elif ext == ".docx":
        text = extract_text_from_docx_bytes(content)
    elif ext in {".xlsx", ".xlsm"}:
        text = extract_text_from_excel_bytes(content)
    elif ext == ".pptx":
        text = extract_text_from_pptx_bytes(content)
    elif ext in {".txt", ".csv"}:
        text = content.decode("utf-8", errors="ignore")
    else:
        text = f"[Unsupported file type: {ext or 'unknown'}]"
    text = re.sub(r"\n{3,}", "\n\n", text or "").strip()
    return {"filename": filename, "extension": ext, "characters": len(text), "text": text}

def compact_doc_text(extracted_files: List[Dict[str, Any]], max_chars: int = 24000) -> str:
    parts = []
    remaining = max_chars
    for f in extracted_files:
        header = f"\n\n===== FILE: {f.get('filename')} =====\n"
        body = f.get("text", "")
        chunk = (header + body)[:remaining]
        parts.append(chunk)
        remaining -= len(chunk)
        if remaining <= 0:
            parts.append("\n[Document text truncated for analysis]\n")
            break
    return "".join(parts).strip()

def evidence_snippets_for_keywords(doc_text: str, keywords: List[str], limit: int = 6) -> List[str]:
    snippets = []
    for sent in split_sentences(doc_text):
        lower = sent.lower()
        if any(k.lower() in lower for k in keywords):
            clean = re.sub(r"\s+", " ", sent).strip()
            if 20 <= len(clean) <= 500:
                snippets.append(clean)
        if len(snippets) >= limit:
            break
    return dedupe_keep_order(snippets, limit)

def infer_doc_sections(ticket: Dict[str, Any], doc_text: str) -> Dict[str, Any]:
    summary = ticket.get("summary") or "Selected Jira ticket"
    description = ticket.get("description") or ""
    scope_evidence = evidence_snippets_for_keywords(doc_text, ["scope", "in scope", "out of scope", "deliver", "capability", "feature", "workflow", "use case"], 8)
    dep_evidence = evidence_snippets_for_keywords(doc_text, ["dependency", "dependencies", "approval", "access", "credential", "environment", "vendor", "security", "firewall", "integration"], 8)
    success_evidence = evidence_snippets_for_keywords(doc_text, ["success", "criteria", "measure", "metric", "target", "threshold", "benefit", "reduction", "improve", "outcome"], 8)
    risk_evidence = evidence_snippets_for_keywords(doc_text, ["risk", "control", "mitigation", "rollback", "constraint", "guardrail", "security", "privacy", "audit"], 8)
    validation_evidence = evidence_snippets_for_keywords(doc_text, ["test", "validation", "verify", "evidence", "sign-off", "acceptance", "done", "demo"], 8)
    stakeholder_evidence = evidence_snippets_for_keywords(doc_text, ["stakeholder", "owner", "team", "approver", "product owner", "technical lead", "business owner"], 6)

    scope = scope_evidence or infer_scope_points(summary, description)
    dependencies = dep_evidence or infer_dependency_points(summary, description)
    success = success_evidence or infer_success_points(summary, description)
    risks = risk_evidence or detect_risk(description).get("evidence") or ["Delivery risks and controls must be confirmed with the relevant stakeholders before implementation."]
    validation = validation_evidence or ["Validation evidence must be captured and attached or referenced on the Jira ticket."]

    acceptance = [
        f"Given the agreed scope and prerequisites are available, when the work is implemented, then {success[0] if success else 'the expected business outcome is achieved'}.",
        f"Given implementation is completed, when validation is performed, then {validation[0]}."
    ]
    acceptance.extend(infer_acceptance_criteria(summary, description)[:2])
    acceptance = dedupe_keep_order(acceptance, 5)

    dor = infer_dor_points(summary, description)
    if dep_evidence:
        dor = dedupe_keep_order(dor + ["Key dependencies from the uploaded supporting documents are confirmed and assigned."], 6)
    if stakeholder_evidence:
        dor = dedupe_keep_order(dor + ["Relevant stakeholder/owner from the supporting documents is identified."], 6)

    dod = infer_dod_points(summary, description)
    if validation_evidence:
        dod = dedupe_keep_order(dod + ["Validation evidence from the agreed test or review process is captured."], 6)

    business_need = evidence_snippets_for_keywords(doc_text, ["business need", "objective", "purpose", "overview", "problem", "benefit"], 3)
    context = evidence_snippets_for_keywords(doc_text, ["overview", "background", "context", "current", "as-is", "problem"], 4)

    return {
        "business_need": first_nonempty(business_need) or f"Deliver {summary} using the supporting evidence provided in the uploaded documents.",
        "context": " ".join(context[:2]) if context else first_nonempty(detect_business_context(description).get("evidence", [])) or "Supporting documents were uploaded and analyzed to strengthen the ticket detail.",
        "problem_statement": first_nonempty(evidence_snippets_for_keywords(doc_text, ["problem", "challenge", "manual", "delay", "risk", "gap"], 3)) or guess_problem_statement(summary, description),
        "scope": scope,
        "dependencies": dependencies,
        "success_criteria": success,
        "acceptance_criteria": acceptance,
        "definition_of_ready": dor,
        "definition_of_done": dod,
        "risks_controls": risks,
        "document_evidence": {
            "scope": scope_evidence, "dependencies": dep_evidence, "success_criteria": success_evidence,
            "risks_controls": risk_evidence, "validation": validation_evidence, "stakeholders": stakeholder_evidence,
        }
    }

def build_description_from_doc_sections(summary: str, sections: Dict[str, Any]) -> str:
    def bullets(items: List[str], max_items: int) -> str:
        return "\n".join(f"- {x}" for x in dedupe_keep_order(items or [], max_items))
    return f"""Business Need:
{sections.get('business_need')}

Context:
{sections.get('context')}

Problem Statement:
{sections.get('problem_statement')}

Scope:
{bullets(sections.get('scope', []), 6)}

Dependencies:
{bullets(sections.get('dependencies', []), 6)}

Success Criteria:
{bullets(sections.get('success_criteria', []), 5)}

Acceptance Criteria:
{bullets(sections.get('acceptance_criteria', []), 5)}

Definition of Ready:
{bullets(sections.get('definition_of_ready', []), 6)}

Definition of Done:
{bullets(sections.get('definition_of_done', []), 6)}

Risks / Controls:
{bullets(sections.get('risks_controls', []), 5)}
""".strip()

def call_openai_for_doc_enhancement(ticket: Dict[str, Any], assessment: Dict[str, Any], doc_text: str, deterministic_sections: Dict[str, Any]) -> Optional[Dict[str, Any]]:
    if not OPENAI_API_KEY:
        return None
    try:
        prompt = f"""You are an expert OSS delivery analyst. Use the Jira ticket and uploaded document evidence to produce compliance-ready ticket content.
Do not invent facts. Return strict JSON with keys: business_need, context, problem_statement, scope, dependencies, success_criteria, acceptance_criteria, definition_of_ready, definition_of_done, risks_controls.
List keys must be arrays of concise strings except business_need, context and problem_statement.
Jira ticket: {ticket}
Current quality assessment: {assessment}
Document-derived draft: {deterministic_sections}
Uploaded document text: {doc_text[:16000]}"""
        response = requests.post(
            "https://api.openai.com/v1/chat/completions",
            headers={"Authorization": f"Bearer {OPENAI_API_KEY}", "Content-Type": "application/json"},
            json={"model": OPENAI_MODEL, "messages": [{"role": "system", "content": "You return only valid JSON. No markdown."}, {"role": "user", "content": prompt}], "temperature": 0.2},
            timeout=60,
        )
        if response.status_code != 200:
            return None
        import json
        return json.loads(response.json()["choices"][0]["message"]["content"])
    except Exception:
        return None

def get_live_ticket_or_error(issue_key: str):
    url = f"{JIRA_BASE_URL}/rest/api/3/issue/{issue_key}"
    params = {"fields": "summary,description,status,assignee,issuetype,labels,priority,components,created,updated,reporter"}
    response = requests.get(url, headers=jira_auth_headers(), params=params, timeout=30)
    if response.status_code != 200:
        return None, JSONResponse(status_code=response.status_code, content=safe_json_response(response))
    return normalize_live_issue(response.json()), None



# -----------------------------------------------------------------------------
# Manager Capacity Workspace helpers
# -----------------------------------------------------------------------------

def parse_project_keys(value: Any) -> List[str]:
    if value is None:
        return []
    raw_items = value if isinstance(value, list) else re.split(r",|\n|;", str(value))
    keys: List[str] = []
    for item in raw_items:
        key = str(item or "").strip().upper()
        if key and re.match(r"^[A-Z][A-Z0-9_]*$", key):
            keys.append(key)
    return dedupe_keep_order(keys, 50)


def get_likely_story_points_field_id(preferred_field: str = "") -> Dict[str, Any]:
    preferred_field = (preferred_field or "").strip()
    if preferred_field:
        return {"field_id": preferred_field, "source": "request"}
    if JIRA_STORY_POINTS_FIELD:
        return {"field_id": JIRA_STORY_POINTS_FIELD, "source": "env"}
    try:
        response = requests.get(f"{JIRA_BASE_URL}/rest/api/3/field", headers=jira_auth_headers(), timeout=30)
        if response.status_code != 200:
            return {"field_id": "", "source": "not_found", "warning": f"Could not discover Story Points field. Jira returned {response.status_code}."}
        fields = response.json() or []
        candidates = []
        for f in fields:
            name = str(f.get("name", ""))
            fid = str(f.get("id", ""))
            lname = name.lower()
            if fid.startswith("customfield_") and (("story" in lname and "point" in lname) or "estimate" in lname):
                candidates.append({"id": fid, "name": name})
        if candidates:
            preferred = next((c for c in candidates if "story" in c["name"].lower() and "point" in c["name"].lower()), candidates[0])
            return {"field_id": preferred["id"], "field_name": preferred["name"], "source": "auto_discovered", "candidates": candidates[:10]}
    except Exception as e:
        return {"field_id": "", "source": "error", "warning": str(e)}
    return {"field_id": "", "source": "not_found", "warning": "No Story Points / estimate field was discovered. Configure JIRA_STORY_POINTS_FIELD or enter the field id manually."}


def story_points_number(value: Any) -> float:
    if value in (None, ""):
        return 0.0
    try:
        return float(value)
    except Exception:
        return 0.0


def normalize_jira_user(user: Dict[str, Any]) -> Dict[str, Any]:
    return {
        "accountId": user.get("accountId") or "",
        "displayName": user.get("displayName") or user.get("name") or "",
        "emailAddress": user.get("emailAddress") or "",
        "active": bool(user.get("active", True)),
        "accountType": user.get("accountType") or "",
        "avatarUrl": ((user.get("avatarUrls") or {}).get("48x48") or (user.get("avatarUrls") or {}).get("32x32") or user.get("avatarUrl") or ""),
    }


def search_jira_users(query: str, max_results: int = 10) -> Dict[str, Any]:
    """Search Jira users so managers do not need to know the exact assignee name.

    Uses user picker first because it is intended for UI search, then falls back to
    user/search. Jira permissions still apply; without Browse users permission Jira
    may only return exact-name matches.
    """
    query = str(query or "").strip()
    if len(query) < 2:
        return {"query": query, "users": [], "message": "Type at least 2 characters."}
    max_results = max(1, min(int(max_results or 10), 50))
    attempts: List[Dict[str, Any]] = []

    picker_url = f"{JIRA_BASE_URL}/rest/api/3/user/picker"
    picker_response = requests.get(picker_url, headers=jira_auth_headers(), params={"query": query, "maxResults": max_results}, timeout=30)
    attempts.append({"endpoint": "/rest/api/3/user/picker", "status_code": picker_response.status_code})
    if picker_response.status_code == 200:
        body = picker_response.json() or {}
        users = body.get("users") or []
        normalized = [normalize_jira_user(u) for u in users]
        normalized = [u for u in normalized if u.get("accountId") or u.get("displayName")]
        if normalized:
            return {"query": query, "users": normalized, "status_code": 200, "method": "picker", "attempts": attempts}

    search_url = f"{JIRA_BASE_URL}/rest/api/3/user/search"
    search_response = requests.get(search_url, headers=jira_auth_headers(), params={"query": query, "maxResults": max_results}, timeout=30)
    attempts.append({"endpoint": "/rest/api/3/user/search", "status_code": search_response.status_code})
    result = safe_json_response(search_response)
    if search_response.status_code != 200:
        result["query"] = query
        result["attempts"] = attempts
        return result
    users = [normalize_jira_user(u) for u in (search_response.json() or [])]
    active_users = [u for u in users if u.get("active")]
    return {"query": query, "users": active_users or users, "status_code": 200, "method": "user_search", "attempts": attempts}


def status_category_name_from_issue(issue: Dict[str, Any]) -> str:
    fields = issue.get("fields") or {}
    status = fields.get("status") or {}
    category = status.get("statusCategory") or {}
    return str(category.get("name") or status.get("name") or "Unknown")



def sanitize_capacity_date(value: Any) -> str:
    """Accept YYYY-MM-DD only for Jira capacity period filters."""
    value = str(value or "").strip()
    if not value:
        return ""
    if not re.match(r"^\d{4}-\d{2}-\d{2}$", value):
        raise ValueError("Capacity period dates must use YYYY-MM-DD format")
    return value

def next_capacity_day(value: str) -> str:
    """Return YYYY-MM-DD for the day after value. Used to make end-date filters inclusive."""
    value = sanitize_capacity_date(value)
    if not value:
        return ""
    return (datetime.strptime(value, "%Y-%m-%d") + timedelta(days=1)).strftime("%Y-%m-%d")



def count_business_days_inclusive(start_date: str = "", end_date: str = "") -> int:
    """Count Monday-Friday days in an inclusive selected capacity period."""
    start_date = sanitize_capacity_date(start_date)
    end_date = sanitize_capacity_date(end_date)
    if not start_date or not end_date:
        return 10  # default two-week sprint assumption
    if start_date > end_date:
        raise ValueError("Capacity start date cannot be after end date")
    start = datetime.strptime(start_date, "%Y-%m-%d")
    end = datetime.strptime(end_date, "%Y-%m-%d")
    days = 0
    current = start
    while current <= end:
        if current.weekday() < 5:
            days += 1
        current += timedelta(days=1)
    return days

def calculated_default_capacity_points(start_date: str = "", end_date: str = "", points_per_working_day: float = 2.0) -> Dict[str, Any]:
    """Return a period-based default capacity while still allowing manager override.

    Assumption: 2 story points per working day. This keeps the old 20-point
    default for a normal 10-working-day sprint, while scaling automatically
    for shorter or longer selected periods.
    """
    try:
        ppd = float(points_per_working_day or 2.0)
    except Exception:
        ppd = 2.0
    if ppd <= 0:
        ppd = 2.0
    working_days = count_business_days_inclusive(start_date, end_date)
    points = round(working_days * ppd, 1)
    return {
        "capacity_points": points,
        "working_days": working_days,
        "points_per_working_day": ppd,
        "basis": f"{working_days} working day(s) × {ppd:g} point(s)/day",
        "can_override": True,
    }

def capacity_period_clause(start_date: str = "", end_date: str = "") -> str:
    """Build the Jira JQL period clause used by capacity checks.

    IMPORTANT: the previous v4 logic allowed active tickets without due dates to
    pass through every period, so the capacity result could remain unchanged even
    when the manager selected different dates.

    This version treats the selected period as a real reporting window. A ticket
    is included only when at least one planning/activity date intersects the
    selected window:
      - created during the period, OR
      - updated during the period, OR
      - due during the period.

    Done/completed tickets are still excluded separately by statusCategory/status.
    """
    start_date = sanitize_capacity_date(start_date)
    end_date = sanitize_capacity_date(end_date)

    if start_date and end_date:
        end_exclusive = next_capacity_day(end_date)
        return (
            "("
            f'(created >= "{start_date}" AND created < "{end_exclusive}") OR '
            f'(updated >= "{start_date}" AND updated < "{end_exclusive}") OR '
            f'(duedate >= "{start_date}" AND duedate <= "{end_date}")'
            ")"
        )

    if start_date:
        return f'(created >= "{start_date}" OR updated >= "{start_date}" OR duedate >= "{start_date}")'

    if end_date:
        end_exclusive = next_capacity_day(end_date)
        return f'(created < "{end_exclusive}" OR updated < "{end_exclusive}" OR duedate <= "{end_date}")'

    return ""

def split_jql_order_by(jql: str):
    raw = str(jql or "").strip()
    match = re.search(r"\s+ORDER\s+BY\s+", raw, flags=re.IGNORECASE)
    if not match:
        return raw, "ORDER BY updated DESC"
    return raw[:match.start()].strip(), raw[match.start():].strip()

def apply_capacity_period_to_jql(base_jql: str, start_date: str = "", end_date: str = "") -> str:
    period = capacity_period_clause(start_date, end_date)
    if not period:
        return base_jql
    body, order_by = split_jql_order_by(base_jql)
    if not body:
        return f"{period} {order_by}"
    return f"({body}) AND {period} {order_by}"


# -----------------------------------------------------------------------------
# Refined capacity algorithm
# -----------------------------------------------------------------------------

def parse_jira_date(value: Any) -> str:
    """Return YYYY-MM-DD from a Jira date/datetime string, or blank."""
    value = str(value or "").strip()
    if not value:
        return ""
    m = re.match(r"^(\d{4}-\d{2}-\d{2})", value)
    return m.group(1) if m else ""


def date_in_capacity_period(date_value: Any, start_date: str = "", end_date: str = "") -> bool:
    d = parse_jira_date(date_value)
    if not d:
        return False
    start_date = sanitize_capacity_date(start_date)
    end_date = sanitize_capacity_date(end_date)
    if start_date and d < start_date:
        return False
    if end_date and d > end_date:
        return False
    return True


def date_before_capacity_start(date_value: Any, start_date: str = "") -> bool:
    d = parse_jira_date(date_value)
    start_date = sanitize_capacity_date(start_date)
    return bool(d and start_date and d < start_date)


def has_capacity_period(start_date: str = "", end_date: str = "") -> bool:
    return bool(sanitize_capacity_date(start_date) or sanitize_capacity_date(end_date))


def capacity_status_from_utilization(utilization: float, capacity_points: float) -> str:
    if capacity_points <= 0:
        return "UNKNOWN"
    if utilization <= 80:
        return "AVAILABLE"
    if utilization <= 110:
        return "NEAR_CAPACITY"
    return "OVER_ALLOCATED"


def calculate_capacity_risk_score(utilization: float, avg_quality: float, unplanned_count: int, unpointed_count: int, overdue_count: int) -> Dict[str, Any]:
    """Combine utilisation + planning quality + Jira ticket quality into one manager risk signal."""
    quality_penalty = max(0.0, 80.0 - float(avg_quality or 0)) * 0.5
    planning_penalty = (float(unplanned_count or 0) * 5.0) + (float(unpointed_count or 0) * 5.0) + (float(overdue_count or 0) * 7.0)
    risk_score = round(float(utilization or 0) + quality_penalty + planning_penalty, 1)
    if risk_score >= 125:
        status = "HIGH_RISK"
    elif risk_score >= 95:
        status = "WATCH"
    else:
        status = "HEALTHY"
    return {
        "capacity_risk_score": risk_score,
        "capacity_risk_status": status,
        "quality_penalty": round(quality_penalty, 1),
        "planning_penalty": round(planning_penalty, 1),
    }


def classify_ticket_for_capacity(ticket: Dict[str, Any], start_date: str = "", end_date: str = "") -> Dict[str, Any]:
    """Classify whether and how a ticket consumes capacity for the selected manager period.

    Principle:
    - Story points are only counted as allocated capacity when the work is planned/touched in the selected period.
    - Active work with no planning date is surfaced as unplanned risk, but does not silently inflate allocated capacity.
    - Done/completed tickets are filtered before this function is called.
    """
    start_date = sanitize_capacity_date(start_date)
    end_date = sanitize_capacity_date(end_date)
    period_selected = has_capacity_period(start_date, end_date)
    status_category = str(ticket.get("statusCategory") or "").strip().lower()
    status_name = str(ticket.get("status") or "").strip().lower()
    story_points = float(ticket.get("storyPoints") or 0)
    due_date = parse_jira_date(ticket.get("dueDate"))
    created = parse_jira_date(ticket.get("created"))
    updated = parse_jira_date(ticket.get("updated"))

    if not period_selected:
        return {
            "includedInCapacity": True,
            "allocatedStoryPoints": story_points,
            "capacityBucket": "ACTIVE_NO_PERIOD",
            "capacityInclusionReason": "No period selected; all active To Do / In Progress work is counted.",
            "unplannedCapacityRisk": False,
            "overdueCapacityRisk": False,
        }

    if due_date and date_in_capacity_period(due_date, start_date, end_date):
        return {
            "includedInCapacity": True,
            "allocatedStoryPoints": story_points,
            "capacityBucket": "PLANNED_DUE_IN_PERIOD",
            "capacityInclusionReason": "Due date falls inside the selected capacity period.",
            "unplannedCapacityRisk": False,
            "overdueCapacityRisk": False,
        }

    if due_date and date_before_capacity_start(due_date, start_date) and status_category == "in progress":
        return {
            "includedInCapacity": True,
            "allocatedStoryPoints": story_points,
            "capacityBucket": "OVERDUE_IN_PROGRESS",
            "capacityInclusionReason": "Ticket is overdue before the selected period and still In Progress; counted as carry-over load.",
            "unplannedCapacityRisk": False,
            "overdueCapacityRisk": True,
        }

    if status_category == "in progress" and date_in_capacity_period(updated, start_date, end_date):
        return {
            "includedInCapacity": True,
            "allocatedStoryPoints": story_points,
            "capacityBucket": "IN_PROGRESS_TOUCHED_IN_PERIOD",
            "capacityInclusionReason": "Ticket is In Progress and was updated during the selected period.",
            "unplannedCapacityRisk": False,
            "overdueCapacityRisk": False,
        }

    if status_category == "to do" and date_in_capacity_period(created, start_date, end_date):
        return {
            "includedInCapacity": True,
            "allocatedStoryPoints": story_points,
            "capacityBucket": "NEW_TODO_CREATED_IN_PERIOD",
            "capacityInclusionReason": "To Do ticket was created during the selected period; counted as newly committed backlog.",
            "unplannedCapacityRisk": False,
            "overdueCapacityRisk": False,
        }

    if not due_date and status_category == "in progress":
        return {
            "includedInCapacity": True,
            "allocatedStoryPoints": 0.0,
            "capacityBucket": "UNPLANNED_IN_PROGRESS_NO_DUE_DATE",
            "capacityInclusionReason": "In Progress ticket has no due/planning date. Shown as unplanned active risk, not counted as planned capacity for this period.",
            "unplannedCapacityRisk": True,
            "overdueCapacityRisk": False,
        }

    if not due_date and status_category == "to do":
        return {
            "includedInCapacity": True,
            "allocatedStoryPoints": 0.0,
            "capacityBucket": "UNPLANNED_TODO_NO_DUE_DATE",
            "capacityInclusionReason": "To Do ticket has no due/planning date. Shown as unplanned backlog risk, not counted as planned capacity for this period.",
            "unplannedCapacityRisk": True,
            "overdueCapacityRisk": False,
        }

    return {
        "includedInCapacity": False,
        "allocatedStoryPoints": 0.0,
        "capacityBucket": "OUTSIDE_SELECTED_PERIOD",
        "capacityInclusionReason": "Active ticket has no planning/activity signal inside the selected period.",
        "unplannedCapacityRisk": False,
        "overdueCapacityRisk": False,
    }




def capacity_period_business_days(start_date: str = "", end_date: str = "") -> List[str]:
    """Return YYYY-MM-DD working days in the selected capacity period."""
    start_date = sanitize_capacity_date(start_date)
    end_date = sanitize_capacity_date(end_date)
    if not start_date or not end_date:
        return []
    if start_date > end_date:
        raise ValueError("Capacity start date cannot be after end date")
    days: List[str] = []
    current = datetime.strptime(start_date, "%Y-%m-%d")
    end = datetime.strptime(end_date, "%Y-%m-%d")
    while current <= end:
        if current.weekday() < 5:
            days.append(current.strftime("%Y-%m-%d"))
        current += timedelta(days=1)
    return days


def clamp_date_to_capacity_period(value: Any, start_date: str, end_date: str) -> str:
    d = parse_jira_date(value)
    start_date = sanitize_capacity_date(start_date)
    end_date = sanitize_capacity_date(end_date)
    if not d:
        return ""
    if start_date and d < start_date:
        return start_date
    if end_date and d > end_date:
        return end_date
    return d


def business_days_between(start_date: str, end_date: str) -> List[str]:
    return capacity_period_business_days(start_date, end_date)


def add_points_evenly(series_map: Dict[str, Dict[str, float]], days: List[str], points: float, key: str = "allocated_points") -> None:
    if not days or not points:
        return
    share = float(points) / len(days)
    for d in days:
        if d in series_map:
            series_map[d][key] = round(series_map[d].get(key, 0.0) + share, 4)


def build_daily_utilization_series(
    tickets: List[Dict[str, Any]],
    start_date: str = "",
    end_date: str = "",
    total_capacity_points: float = 0,
) -> Dict[str, Any]:
    """Build a day-by-day utilization view for the selected period.

    Jira Story Points are not daily worklog hours, so this is a planning model:
      - due-date work is spread from the period start to its due date;
      - overdue and in-progress work is spread across the full selected period;
      - newly created To Do work is spread from created date to period end;
      - unplanned/no-date work stays visible as risk but contributes 0 planned daily load.
    """
    days = capacity_period_business_days(start_date, end_date)
    if not days:
        return {
            "available": False,
            "message": "Select both start and end dates to render daily utilization.",
            "days": [],
            "daily_capacity_points": 0,
            "series": [],
        }

    capacity = float(total_capacity_points or 0)
    daily_capacity = round(capacity / len(days), 2) if capacity > 0 and days else 0
    series_map: Dict[str, Dict[str, float]] = {
        d: {
            "date": d,
            "daily_capacity_points": daily_capacity,
            "allocated_points": 0.0,
            "raw_story_points_due_or_touched": 0.0,
            "ticket_count_signal": 0,
            "utilization_percent": 0.0,
        }
        for d in days
    }

    for t in tickets or []:
        allocated = float(t.get("allocatedStoryPoints") or t.get("allocated_story_points") or 0)
        raw_sp = float(t.get("storyPoints") or t.get("raw_story_points") or 0)
        if allocated <= 0:
            continue
        bucket = str(t.get("capacityBucket") or "").upper()
        status_category = str(t.get("statusCategory") or "").lower()
        due_date = clamp_date_to_capacity_period(t.get("dueDate"), start_date, end_date)
        created = clamp_date_to_capacity_period(t.get("created"), start_date, end_date)
        updated = clamp_date_to_capacity_period(t.get("updated"), start_date, end_date)

        allocation_days: List[str] = []
        if bucket == "PLANNED_DUE_IN_PERIOD" and due_date:
            allocation_days = business_days_between(days[0], due_date)
        elif bucket == "NEW_TODO_CREATED_IN_PERIOD" and created:
            allocation_days = business_days_between(created, days[-1])
        elif bucket in {"OVERDUE_IN_PROGRESS", "IN_PROGRESS_TOUCHED_IN_PERIOD"} or status_category == "in progress":
            allocation_days = days[:]
        elif updated:
            allocation_days = business_days_between(updated, days[-1])
        else:
            allocation_days = days[:]

        allocation_days = [d for d in allocation_days if d in series_map]
        if not allocation_days:
            allocation_days = days[:]
        add_points_evenly(series_map, allocation_days, allocated, "allocated_points")

        signal_date = due_date or updated or created
        if signal_date in series_map:
            series_map[signal_date]["raw_story_points_due_or_touched"] = round(series_map[signal_date].get("raw_story_points_due_or_touched", 0.0) + raw_sp, 4)
            series_map[signal_date]["ticket_count_signal"] = int(series_map[signal_date].get("ticket_count_signal", 0)) + 1

    series = []
    for d in days:
        row = series_map[d]
        row["allocated_points"] = round(row.get("allocated_points", 0.0), 2)
        row["raw_story_points_due_or_touched"] = round(row.get("raw_story_points_due_or_touched", 0.0), 2)
        row["utilization_percent"] = round((row["allocated_points"] / daily_capacity) * 100, 1) if daily_capacity > 0 else 0
        series.append(row)

    peak = max((r["utilization_percent"] for r in series), default=0)
    return {
        "available": True,
        "model": "allocated_story_points_spread_across_business_days",
        "basis": "Daily capacity = selected period capacity divided by working days. Allocated points are spread using capacity bucket logic.",
        "working_days": len(days),
        "daily_capacity_points": daily_capacity,
        "peak_utilization_percent": round(peak, 1),
        "over_capacity_days": sum(1 for r in series if r["utilization_percent"] > 100),
        "days": days,
        "series": series,
    }


def active_capacity_base_jql(project_keys: List[str], assignee_clause: str) -> str:
    project_clause = ""
    if project_keys:
        project_clause = "project in (" + ", ".join(project_keys) + ") AND "
    return f"{project_clause}{assignee_clause} AND statusCategory != Done ORDER BY updated DESC"


def ticket_from_jira_issue(issue: Dict[str, Any], sp_field: str, default_assignee: str = "", default_account_id: str = "") -> Dict[str, Any]:
    f = issue.get("fields") or {}
    status = f.get("status") or {}
    category = status_category_name_from_issue(issue)
    key = issue.get("key") or ""
    points = story_points_number(f.get(sp_field)) if sp_field else 0.0
    assignee_obj = f.get("assignee") or {}
    assignee_name = (assignee_obj.get("displayName") if isinstance(assignee_obj, dict) else "") or default_assignee or "Unassigned"
    assignee_account_id = (assignee_obj.get("accountId") if isinstance(assignee_obj, dict) else "") or default_account_id
    return {
        "issueKey": key,
        "projectKey": key.split("-")[0] if "-" in key else "",
        "summary": f.get("summary") or "",
        "issueType": (f.get("issuetype") or {}).get("name", ""),
        "status": status.get("name") or "",
        "statusCategory": category,
        "priority": (f.get("priority") or {}).get("name", ""),
        "storyPoints": points,
        "hasStoryPoints": bool(points),
        "created": f.get("created") or "",
        "updated": f.get("updated") or "",
        "dueDate": f.get("duedate") or "",
        "parent": ((f.get("parent") or {}).get("key") if isinstance(f.get("parent"), dict) else "") or "",
        "assignee": assignee_name,
        "assigneeAccountId": assignee_account_id,
    }


def fetch_capacity_issues(jql: str, fields: List[str], max_results: int = 0) -> Dict[str, Any]:
    url = f"{JIRA_BASE_URL}/rest/api/3/search/jql"
    page_size = 100
    remaining = None if not max_results or int(max_results) <= 0 else int(max_results)
    next_page_token = None
    issues_out: List[Dict[str, Any]] = []
    fetched_count = 0
    page_count = 0
    while True:
        request_size = page_size if remaining is None else max(1, min(page_size, remaining))
        params = {"jql": jql, "maxResults": request_size, "fields": ",".join(fields), "fieldsByKeys": "false"}
        if next_page_token:
            params["nextPageToken"] = next_page_token
        response = requests.get(url, headers=jira_auth_headers(), params=params, timeout=120)
        if response.status_code != 200:
            raise RuntimeError(str(safe_json_response(response)))
        body = response.json() or {}
        page_count += 1
        issues = body.get("issues", []) or []
        fetched_count += len(issues)
        issues_out.extend(issues)
        if remaining is not None:
            remaining -= len(issues)
            if remaining <= 0:
                break
        next_page_token = body.get("nextPageToken")
        if body.get("isLast", False) or not next_page_token or not issues:
            break
    return {"issues": issues_out, "jira_fetched_count": fetched_count, "jira_pages_read": page_count}


def build_capacity_dashboard(
    assignee: str,
    project_keys: List[str],
    story_points_field_id: str = "",
    sprint_capacity_points: float = 20,
    jql_override: str = "",
    max_results: int = 0,
    assignee_account_id: str = "",
    start_date: str = "",
    end_date: str = "",
) -> Dict[str, Any]:
    assignee = str(assignee or "").strip()
    assignee_account_id = str(assignee_account_id or "").strip()
    start_date = sanitize_capacity_date(start_date)
    end_date = sanitize_capacity_date(end_date)
    if start_date and end_date and start_date > end_date:
        raise ValueError("Capacity start date cannot be after end date")
    default_capacity_calc = calculated_default_capacity_points(start_date, end_date)
    if not sprint_capacity_points or float(sprint_capacity_points or 0) <= 0:
        sprint_capacity_points = default_capacity_calc["capacity_points"]
    if not assignee and not assignee_account_id:
        raise ValueError("assignee or assigneeAccountId is required")

    sp_info = get_likely_story_points_field_id(story_points_field_id)
    sp_field = sp_info.get("field_id") or ""

    if jql_override and str(jql_override).strip():
        jql = str(jql_override).strip()
        assignee_clause_used = "jql_override"
    else:
        if assignee_account_id:
            assignee_clause = f"assignee = {quote_jql_value(assignee_account_id)}"
            assignee_clause_used = "accountId"
        else:
            assignee_clause = f"assignee = {quote_jql_value(assignee)}"
            assignee_clause_used = "displayName"
        jql = active_capacity_base_jql(project_keys, assignee_clause)

    fields = ["summary", "description", "status", "assignee", "issuetype", "labels", "priority", "components", "created", "updated", "duedate", "parent", "reporter"]
    if sp_field:
        fields.append(sp_field)

    raw = fetch_capacity_issues(jql, fields, max_results)
    tickets: List[Dict[str, Any]] = []
    excluded_outside_period = 0
    for issue in raw["issues"]:
        f = issue.get("fields") or {}
        status = f.get("status") or {}
        category = status_category_name_from_issue(issue)
        if is_done_or_completed_status(status.get("name")) or str(category).lower() == "done":
            continue
        ticket = ticket_from_jira_issue(issue, sp_field, assignee, assignee_account_id)
        classification = classify_ticket_for_capacity(ticket, start_date, end_date)
        ticket.update(classification)
        if not classification.get("includedInCapacity"):
            excluded_outside_period += 1
            continue
        tickets.append(ticket)

    allocated_points = round(sum(float(t.get("allocatedStoryPoints") or 0) for t in tickets), 1)
    raw_points = round(sum(float(t.get("storyPoints") or 0) for t in tickets), 1)
    todo_points = round(sum(float(t.get("allocatedStoryPoints") or 0) for t in tickets if str(t.get("statusCategory", "")).lower() == "to do"), 1)
    in_progress_points = round(sum(float(t.get("allocatedStoryPoints") or 0) for t in tickets if str(t.get("statusCategory", "")).lower() == "in progress"), 1)
    unpointed_count = sum(1 for t in tickets if not t.get("hasStoryPoints"))
    unplanned_count = sum(1 for t in tickets if t.get("unplannedCapacityRisk"))
    overdue_count = sum(1 for t in tickets if t.get("overdueCapacityRisk"))

    capacity = float(sprint_capacity_points or 0)
    remaining_capacity = round(capacity - allocated_points, 1) if capacity > 0 else 0
    utilization = round((allocated_points / capacity) * 100, 1) if capacity > 0 else 0
    capacity_status = capacity_status_from_utilization(utilization, capacity)
    risk = calculate_capacity_risk_score(utilization, 100, unplanned_count, unpointed_count, overdue_count)

    by_project: Dict[str, Dict[str, Any]] = {}
    by_status_category: Dict[str, Dict[str, Any]] = {}
    by_capacity_bucket: Dict[str, Dict[str, Any]] = {}
    for t in tickets:
        allocated = float(t.get("allocatedStoryPoints") or 0)
        raw_sp = float(t.get("storyPoints") or 0)
        for key, target in [(t["projectKey"] or "UNKNOWN", by_project), (t["statusCategory"] or "Unknown", by_status_category), (t.get("capacityBucket") or "UNKNOWN", by_capacity_bucket)]:
            if key not in target:
                target[key] = {"name": key, "ticket_count": 0, "story_points": 0.0, "raw_story_points": 0.0, "unpointed": 0}
            target[key]["ticket_count"] += 1
            target[key]["story_points"] = round(target[key]["story_points"] + allocated, 1)
            target[key]["raw_story_points"] = round(target[key]["raw_story_points"] + raw_sp, 1)
            target[key]["unpointed"] += 0 if t.get("hasStoryPoints") else 1

    tickets_sorted = sorted(tickets, key=lambda x: (-float(x.get("allocatedStoryPoints") or 0), -float(x.get("storyPoints") or 0), x.get("projectKey") or "", x.get("issueKey") or ""))
    return {
        "assignee": assignee or (tickets_sorted[0]["assignee"] if tickets_sorted else ""),
        "assigneeAccountId": assignee_account_id or (tickets_sorted[0]["assigneeAccountId"] if tickets_sorted else ""),
        "projectKeys": project_keys,
        "jql": jql,
        "assigneeClauseUsed": assignee_clause_used,
        "storyPointsField": sp_info,
        "sprintCapacityPoints": capacity,
        "capacityDefault": default_capacity_calc,
        "capacityPeriod": {
            "startDate": start_date,
            "endDate": end_date,
            "periodClause": capacity_period_clause(start_date, end_date),
            "filterMode": "Refined algorithm: active Jira tickets are fetched first; backend classifies each ticket by due date, carry-over overdue work, in-period updates/creation, and unplanned no-date risk.",
            "algorithm": "planned_due > overdue_in_progress > in_progress_updated > todo_created > unplanned_no_due_risk > outside_period",
        },
        "summary": {
            "active_ticket_count": len(tickets_sorted),
            "allocated_story_points": allocated_points,
            "total_story_points": allocated_points,
            "raw_story_points": raw_points,
            "todo_story_points": todo_points,
            "in_progress_story_points": in_progress_points,
            "unpointed_ticket_count": unpointed_count,
            "unplanned_ticket_count": unplanned_count,
            "overdue_ticket_count": overdue_count,
            "excluded_outside_period_count": excluded_outside_period,
            "remaining_capacity_points": remaining_capacity,
            "utilization_percent": utilization,
            "capacity_status": capacity_status,
            **risk,
        },
        "by_project": sorted(by_project.values(), key=lambda x: (-x["story_points"], x["name"])),
        "by_status_category": sorted(by_status_category.values(), key=lambda x: x["name"]),
        "by_capacity_bucket": sorted(by_capacity_bucket.values(), key=lambda x: (-x["story_points"], x["name"])),
        "daily_utilization": build_daily_utilization_series(tickets_sorted, start_date, end_date, capacity),
        "tickets": tickets_sorted,
        "jira_fetched_count": raw["jira_fetched_count"],
        "jira_pages_read": raw["jira_pages_read"],
    }


def build_team_capacity_dashboard(
    members: List[Dict[str, Any]],
    project_keys: List[str],
    story_points_field_id: str = "",
    default_capacity_points: float = 20,
    jql_override: str = "",
    max_results: int = 0,
    start_date: str = "",
    end_date: str = "",
) -> Dict[str, Any]:
    start_date = sanitize_capacity_date(start_date)
    end_date = sanitize_capacity_date(end_date)
    if start_date and end_date and start_date > end_date:
        raise ValueError("Capacity start date cannot be after end date")
    default_capacity_calc = calculated_default_capacity_points(start_date, end_date)
    if not default_capacity_points or float(default_capacity_points or 0) <= 0:
        default_capacity_points = default_capacity_calc["capacity_points"]
    clean_members: List[Dict[str, Any]] = []
    seen = set()
    for m in members or []:
        name = str(m.get("displayName") or m.get("assignee") or m.get("name") or "").strip()
        account_id = str(m.get("accountId") or m.get("assigneeAccountId") or "").strip()
        try:
            capacity = float(m.get("capacityPoints") if m.get("capacityPoints") not in (None, "") else default_capacity_points)
        except Exception:
            capacity = float(default_capacity_points or 0)
        key = account_id or name.lower()
        if not key or key in seen:
            continue
        seen.add(key)
        clean_members.append({"displayName": name or account_id, "accountId": account_id, "capacityPoints": capacity})
    if not clean_members:
        raise ValueError("At least one team member is required")

    sp_info = get_likely_story_points_field_id(story_points_field_id)
    sp_field = sp_info.get("field_id") or ""

    if jql_override and str(jql_override).strip():
        jql = str(jql_override).strip()
        member_clause_used = "jql_override"
    else:
        assignee_values = []
        for m in clean_members:
            assignee_values.append(quote_jql_value(m.get("accountId") or m.get("displayName")))
        assignee_clause = "assignee in (" + ", ".join(assignee_values) + ")"
        member_clause_used = "accountId/displayName roster"
        jql = active_capacity_base_jql(project_keys, assignee_clause)

    fields = ["summary", "description", "status", "assignee", "issuetype", "labels", "priority", "components", "created", "updated", "duedate", "parent", "reporter"]
    if sp_field:
        fields.append(sp_field)

    member_lookup: Dict[str, Dict[str, Any]] = {}
    for m in clean_members:
        if m.get("accountId"):
            member_lookup[m.get("accountId")] = m
        if m.get("displayName"):
            member_lookup[str(m.get("displayName")).lower()] = m

    raw = fetch_capacity_issues(jql, fields, max_results)
    tickets: List[Dict[str, Any]] = []
    excluded_outside_period = 0
    for issue in raw["issues"]:
        f = issue.get("fields") or {}
        status = f.get("status") or {}
        category = status_category_name_from_issue(issue)
        if is_done_or_completed_status(status.get("name")) or str(category).lower() == "done":
            continue
        base_ticket = ticket_from_jira_issue(issue, sp_field)
        classification = classify_ticket_for_capacity(base_ticket, start_date, end_date)
        base_ticket.update(classification)
        if not classification.get("includedInCapacity"):
            excluded_outside_period += 1
            continue
        member = member_lookup.get(base_ticket.get("assigneeAccountId")) or member_lookup.get(str(base_ticket.get("assignee") or "").lower()) or {"displayName": base_ticket.get("assignee"), "accountId": base_ticket.get("assigneeAccountId"), "capacityPoints": default_capacity_points}
        normalized = normalize_live_issue(issue)
        assessment = evaluate_ticket(normalized)
        base_ticket.update({
            "memberCapacityPoints": float(member.get("capacityPoints") or 0),
            "qualityScore": assessment.get("score", 0),
            "qualityStatus": assessment.get("status", ""),
            "qualityCompliance": assessment.get("compliance", ""),
            "qualityIssues": assessment.get("issues", [])[:4],
        })
        tickets.append(base_ticket)

    by_member: Dict[str, Dict[str, Any]] = {}
    for m in clean_members:
        key = m.get("accountId") or str(m.get("displayName", "")).lower()
        by_member[key] = {
            "displayName": m.get("displayName"),
            "accountId": m.get("accountId"),
            "capacityPoints": float(m.get("capacityPoints") or 0),
            "active_ticket_count": 0,
            "allocated_story_points": 0.0,
            "total_story_points": 0.0,
            "raw_story_points": 0.0,
            "todo_story_points": 0.0,
            "in_progress_story_points": 0.0,
            "unpointed_ticket_count": 0,
            "unplanned_ticket_count": 0,
            "overdue_ticket_count": 0,
            "quality_score_total": 0.0,
            "quality_ticket_count": 0,
            "average_quality_score": 0,
            "pass_count": 0,
            "partial_count": 0,
            "fail_count": 0,
            "utilization_percent": 0,
            "remaining_capacity_points": float(m.get("capacityPoints") or 0),
            "capacity_status": "AVAILABLE" if float(m.get("capacityPoints") or 0) > 0 else "UNKNOWN",
            "capacity_risk_score": 0,
            "capacity_risk_status": "UNKNOWN",
        }

    by_project: Dict[str, Dict[str, Any]] = {}
    by_quality_gap: Dict[str, int] = {}
    by_capacity_bucket: Dict[str, Dict[str, Any]] = {}
    for t in tickets:
        key = t.get("assigneeAccountId") or str(t.get("assignee", "")).lower()
        row = by_member.get(key)
        if row is None:
            row = by_member.setdefault(key or t.get("assignee"), {
                "displayName": t.get("assignee") or "Unassigned", "accountId": t.get("assigneeAccountId") or "", "capacityPoints": float(default_capacity_points or 0),
                "active_ticket_count": 0, "allocated_story_points": 0.0, "total_story_points": 0.0, "raw_story_points": 0.0, "todo_story_points": 0.0, "in_progress_story_points": 0.0,
                "unpointed_ticket_count": 0, "unplanned_ticket_count": 0, "overdue_ticket_count": 0, "quality_score_total": 0.0, "quality_ticket_count": 0,
                "average_quality_score": 0, "pass_count": 0, "partial_count": 0, "fail_count": 0, "utilization_percent": 0, "remaining_capacity_points": float(default_capacity_points or 0),
                "capacity_status": "UNKNOWN", "capacity_risk_score": 0, "capacity_risk_status": "UNKNOWN",
            })
        pts = float(t.get("allocatedStoryPoints") or 0)
        raw_pts = float(t.get("storyPoints") or 0)
        row["active_ticket_count"] += 1
        row["allocated_story_points"] = round(row["allocated_story_points"] + pts, 1)
        row["total_story_points"] = row["allocated_story_points"]
        row["raw_story_points"] = round(row["raw_story_points"] + raw_pts, 1)
        if str(t.get("statusCategory", "")).lower() == "to do":
            row["todo_story_points"] = round(row["todo_story_points"] + pts, 1)
        if str(t.get("statusCategory", "")).lower() == "in progress":
            row["in_progress_story_points"] = round(row["in_progress_story_points"] + pts, 1)
        row["unpointed_ticket_count"] += 0 if t.get("hasStoryPoints") else 1
        row["unplanned_ticket_count"] += 1 if t.get("unplannedCapacityRisk") else 0
        row["overdue_ticket_count"] += 1 if t.get("overdueCapacityRisk") else 0
        row["quality_score_total"] += float(t.get("qualityScore") or 0)
        row["quality_ticket_count"] += 1
        qs = t.get("qualityStatus")
        row["pass_count"] += 1 if qs == "PASS" else 0
        row["partial_count"] += 1 if qs == "PARTIAL" else 0
        row["fail_count"] += 1 if qs == "FAIL" else 0
        for issue_text in t.get("qualityIssues") or []:
            by_quality_gap[issue_text] = by_quality_gap.get(issue_text, 0) + 1
        pk = t.get("projectKey") or "UNKNOWN"
        if pk not in by_project:
            by_project[pk] = {"name": pk, "ticket_count": 0, "story_points": 0.0, "raw_story_points": 0.0, "quality_score_total": 0.0, "quality_ticket_count": 0, "average_quality_score": 0}
        by_project[pk]["ticket_count"] += 1
        by_project[pk]["story_points"] = round(by_project[pk]["story_points"] + pts, 1)
        by_project[pk]["raw_story_points"] = round(by_project[pk]["raw_story_points"] + raw_pts, 1)
        by_project[pk]["quality_score_total"] += float(t.get("qualityScore") or 0)
        by_project[pk]["quality_ticket_count"] += 1
        bucket = t.get("capacityBucket") or "UNKNOWN"
        if bucket not in by_capacity_bucket:
            by_capacity_bucket[bucket] = {"name": bucket, "ticket_count": 0, "story_points": 0.0, "raw_story_points": 0.0, "unpointed": 0}
        by_capacity_bucket[bucket]["ticket_count"] += 1
        by_capacity_bucket[bucket]["story_points"] = round(by_capacity_bucket[bucket]["story_points"] + pts, 1)
        by_capacity_bucket[bucket]["raw_story_points"] = round(by_capacity_bucket[bucket]["raw_story_points"] + raw_pts, 1)
        by_capacity_bucket[bucket]["unpointed"] += 0 if t.get("hasStoryPoints") else 1

    for row in by_member.values():
        cap = float(row.get("capacityPoints") or 0)
        allocated = float(row.get("allocated_story_points") or 0)
        row["utilization_percent"] = round((allocated / cap) * 100, 1) if cap > 0 else 0
        row["remaining_capacity_points"] = round(cap - allocated, 1) if cap > 0 else 0
        row["capacity_status"] = capacity_status_from_utilization(row["utilization_percent"], cap)
        row["average_quality_score"] = round(row["quality_score_total"] / max(1, row["quality_ticket_count"]), 1)
        row.update(calculate_capacity_risk_score(row["utilization_percent"], row["average_quality_score"], row["unplanned_ticket_count"], row["unpointed_ticket_count"], row["overdue_ticket_count"]))
        row.pop("quality_score_total", None)
        row.pop("quality_ticket_count", None)
    for row in by_project.values():
        row["average_quality_score"] = round(row["quality_score_total"] / max(1, row["quality_ticket_count"]), 1)
        row.pop("quality_score_total", None)
        row.pop("quality_ticket_count", None)

    total_capacity = round(sum(float(m.get("capacityPoints") or 0) for m in clean_members), 1)
    allocated_points = round(sum(float(t.get("allocatedStoryPoints") or 0) for t in tickets), 1)
    raw_points = round(sum(float(t.get("storyPoints") or 0) for t in tickets), 1)
    quality_scores = [float(t.get("qualityScore") or 0) for t in tickets]
    avg_quality = round(sum(quality_scores) / max(1, len(quality_scores)), 1)
    pass_count = sum(1 for t in tickets if t.get("qualityStatus") == "PASS")
    partial_count = sum(1 for t in tickets if t.get("qualityStatus") == "PARTIAL")
    fail_count = sum(1 for t in tickets if t.get("qualityStatus") == "FAIL")
    utilization = round((allocated_points / total_capacity) * 100, 1) if total_capacity > 0 else 0
    unplanned_count = sum(1 for t in tickets if t.get("unplannedCapacityRisk"))
    unpointed_count = sum(1 for t in tickets if not t.get("hasStoryPoints"))
    overdue_count = sum(1 for t in tickets if t.get("overdueCapacityRisk"))
    team_risk = calculate_capacity_risk_score(utilization, avg_quality, unplanned_count, unpointed_count, overdue_count)

    return {
        "members": clean_members,
        "projectKeys": project_keys,
        "jql": jql,
        "memberClauseUsed": member_clause_used,
        "storyPointsField": sp_info,
        "capacityDefault": default_capacity_calc,
        "capacityPeriod": {
            "startDate": start_date,
            "endDate": end_date,
            "periodClause": capacity_period_clause(start_date, end_date),
            "filterMode": "Refined algorithm: active Jira tickets are fetched first; backend classifies each ticket by due date, carry-over overdue work, in-period updates/creation, and unplanned no-date risk.",
            "algorithm": "planned_due > overdue_in_progress > in_progress_updated > todo_created > unplanned_no_due_risk > outside_period",
        },
        "summary": {
            "team_member_count": len(clean_members),
            "active_ticket_count": len(tickets),
            "total_capacity_points": total_capacity,
            "allocated_story_points": allocated_points,
            "total_story_points": allocated_points,
            "raw_story_points": raw_points,
            "remaining_capacity_points": round(total_capacity - allocated_points, 1) if total_capacity > 0 else 0,
            "utilization_percent": utilization,
            "capacity_status": capacity_status_from_utilization(utilization, total_capacity),
            "unpointed_ticket_count": unpointed_count,
            "unplanned_ticket_count": unplanned_count,
            "overdue_ticket_count": overdue_count,
            "excluded_outside_period_count": excluded_outside_period,
            "average_quality_score": avg_quality,
            "pass_count": pass_count,
            "partial_count": partial_count,
            "fail_count": fail_count,
            "over_allocated_member_count": sum(1 for r in by_member.values() if r.get("capacity_status") == "OVER_ALLOCATED"),
            "available_member_count": sum(1 for r in by_member.values() if r.get("capacity_status") == "AVAILABLE"),
            **team_risk,
        },
        "by_member": sorted(by_member.values(), key=lambda x: (-x.get("capacity_risk_score", 0), -x.get("utilization_percent", 0), x.get("displayName") or "")),
        "by_project": sorted(by_project.values(), key=lambda x: (-x.get("story_points", 0), x.get("name") or "")),
        "by_capacity_bucket": sorted(by_capacity_bucket.values(), key=lambda x: (-x.get("story_points", 0), x.get("name") or "")),
        "quality_gaps": sorted([{"issue": k, "count": v} for k, v in by_quality_gap.items()], key=lambda x: -x["count"]),
        "daily_utilization": build_daily_utilization_series(tickets, start_date, end_date, total_capacity),
        "tickets": sorted(tickets, key=lambda x: (-float(x.get("allocatedStoryPoints") or 0), -float(x.get("storyPoints") or 0), x.get("assignee") or "", x.get("issueKey") or "")),
        "jira_fetched_count": raw["jira_fetched_count"],
        "jira_pages_read": raw["jira_pages_read"],
    }

@app.get("/", response_class=HTMLResponse)

def home(request: Request):
    return templates.TemplateResponse(
        "index.html",
        {
            "request": request,
            "default_project": str(DEFAULT_PROJECT),
            "story_points_field": str(JIRA_STORY_POINTS_FIELD)
        }
    )


@app.get("/manager", response_class=HTMLResponse)
def manager_page(request: Request):
    return templates.TemplateResponse("manager.html", {"request": request, "default_project": DEFAULT_PROJECT, "story_points_field": JIRA_STORY_POINTS_FIELD})


@app.get("/api/health")
def api_health():
    return {"ok": True, "jira_base_url": JIRA_BASE_URL, "default_project": DEFAULT_PROJECT}


@app.get("/api/jira/me")
def jira_me():
    url = f"{JIRA_BASE_URL}/rest/api/3/myself"
    response = requests.get(url, headers=jira_auth_headers(), timeout=30)
    return safe_json_response(response)


@app.get("/api/jira/users/search")
def jira_users_search(query: str = Query(default=""), maxResults: int = Query(default=10)):
    result = search_jira_users(query, maxResults)
    if result.get("status_code") and result.get("status_code") != 200:
        return JSONResponse(status_code=int(result.get("status_code", 500)), content=result)
    return result


@app.get("/api/jira/fields/story-points")
def jira_story_points_fields():
    """Return likely Story Points / estimate custom fields for this Jira site.

    Use the returned `id` value in .env as JIRA_STORY_POINTS_FIELD, e.g. customfield_10016.
    """
    url = f"{JIRA_BASE_URL}/rest/api/3/field"
    response = requests.get(url, headers=jira_auth_headers(), timeout=30)
    result = safe_json_response(response)

    if response.status_code != 200:
        return JSONResponse(status_code=response.status_code, content=result)

    fields = result.get("jira_response", [])
    likely = []
    keywords = ["story point", "story points", "story point estimate", "estimate", "estimation"]
    for f in fields:
        name = str(f.get("name", ""))
        fid = str(f.get("id", ""))
        lname = name.lower()
        if fid.startswith("customfield_") and any(k in lname for k in keywords):
            likely.append({
                "id": fid,
                "name": name,
                "custom": f.get("custom"),
                "schema": f.get("schema", {}),
            })

    return {
        "configured_story_points_field": JIRA_STORY_POINTS_FIELD or None,
        "likely_story_point_fields": likely,
        "instruction": "Set JIRA_STORY_POINTS_FIELD in .env to the correct id, e.g. customfield_10016, then restart uvicorn.",
    }


@app.post("/api/jira/search")
def jira_search(body: Dict[str, Any]):
    project_key = normalize_project_key(body.get("projectKey"))
    jql = body.get("jql") or f"project = {project_key} ORDER BY created DESC"
    max_results = int(body.get("maxResults", 10))
    next_page_token = body.get("nextPageToken")
    fields = body.get("fields", ["summary", "description", "status", "assignee", "issuetype", "labels", "priority", "components", "created", "updated", "reporter"])
    url = f"{JIRA_BASE_URL}/rest/api/3/search/jql"
    params = {"jql": jql, "maxResults": max_results, "fields": ",".join(fields), "fieldsByKeys": "false"}
    if next_page_token:
        params["nextPageToken"] = next_page_token
    response = requests.get(url, headers=jira_auth_headers(), params=params, timeout=60)
    result = safe_json_response(response)
    result["projectKey"] = project_key
    result["jql"] = jql
    return result


@app.get("/api/jira/issue/{issue_key}")
def jira_get_issue(issue_key: str):
    url = f"{JIRA_BASE_URL}/rest/api/3/issue/{issue_key}"
    params = {"fields": "summary,description,status,assignee,issuetype,labels,priority,components,created,updated,reporter,parent"}
    response = requests.get(url, headers=jira_auth_headers(), params=params, timeout=30)
    return safe_json_response(response)


@app.post("/api/tickets/check")
def check_ticket(payload: Dict[str, Any]):
    data = {
        "summary": payload.get("summary", ""),
        "description": payload.get("description", ""),
        "issue_type": payload.get("issue_type") or payload.get("issueType") or "",
        "labels": payload.get("labels", []),
        "components": payload.get("components", []),
        "priority": payload.get("priority", ""),
        "dependencies": payload.get("dependencies", []),
        "acceptance_criteria": payload.get("acceptance_criteria", []),
    }
    return evaluate_ticket(data)


@app.post("/api/tickets/check-live")
def check_live_ticket(payload: Dict[str, Any]):
    issue_key = (payload.get("issueKey") or "").strip().upper()
    if not issue_key:
        return JSONResponse(status_code=400, content={"error": "issueKey is required"})
    url = f"{JIRA_BASE_URL}/rest/api/3/issue/{issue_key}"
    params = {"fields": "summary,description,status,assignee,issuetype,labels,priority,components,created,updated,reporter,parent"}
    response = requests.get(url, headers=jira_auth_headers(), params=params, timeout=30)
    if response.status_code != 200:
        return JSONResponse(status_code=response.status_code, content=safe_json_response(response))
    issue = response.json()
    normalized = normalize_live_issue(issue)
    assessment = evaluate_ticket(normalized)
    return {"issueKey": issue_key, "normalized_issue": normalized, "assessment": assessment}


@app.post("/api/tickets/rewrite")
def rewrite_ticket(payload: Dict[str, Any]):
    return generate_smart_rewrite((payload.get("summary") or "").strip(), (payload.get("description") or "").strip(), (payload.get("issueType") or payload.get("issue_type") or "Story").strip())


@app.post("/api/tickets/rewrite-live")
def rewrite_live_ticket(payload: Dict[str, Any]):
    issue_key = (payload.get("issueKey") or "").strip().upper()
    if not issue_key:
        return JSONResponse(status_code=400, content={"error": "issueKey is required"})
    url = f"{JIRA_BASE_URL}/rest/api/3/issue/{issue_key}"
    params = {"fields": "summary,description,status,assignee,issuetype,labels,priority,components,created,updated,reporter,parent"}
    response = requests.get(url, headers=jira_auth_headers(), params=params, timeout=30)
    if response.status_code != 200:
        return JSONResponse(status_code=response.status_code, content=safe_json_response(response))
    issue = response.json()
    normalized = normalize_live_issue(issue)
    rewrite = generate_smart_rewrite(normalized["summary"], normalized["description"], normalized["issue_type"] or "Story")
    return {"issueKey": issue_key, "normalized_issue": normalized, "rewrite": rewrite}


@app.post("/api/tickets/enhance-with-input")
def enhance_with_human_input(payload: Dict[str, Any]):
    issue_key = (payload.get("issueKey") or "").strip().upper()
    human_input = payload.get("human_input") or {}
    if not issue_key:
        return JSONResponse(status_code=400, content={"error": "issueKey is required"})
    url = f"{JIRA_BASE_URL}/rest/api/3/issue/{issue_key}"
    params = {"fields": "summary,description,status,assignee,issuetype,labels,priority,components,created,updated,reporter,parent"}
    response = requests.get(url, headers=jira_auth_headers(), params=params, timeout=30)
    if response.status_code != 200:
        return JSONResponse(status_code=response.status_code, content=safe_json_response(response))
    issue = response.json()
    normalized = normalize_live_issue(issue)
    enhanced = merge_human_input_into_rewrite(normalized["summary"], normalized["description"], normalized["issue_type"] or "Story", human_input)
    return {"issueKey": issue_key, "normalized_issue": normalized, "human_input": human_input, "enhanced_rewrite": enhanced}


@app.post("/api/tickets/auto-clean")
def auto_clean_ticket(payload: Dict[str, Any]):
    return auto_clean_and_normalize_ticket((payload.get("summary") or "").strip(), (payload.get("description") or "").strip(), (payload.get("issueType") or payload.get("issue_type") or "Story").strip())


@app.post("/api/tickets/auto-clean-live")
def auto_clean_live_ticket(payload: Dict[str, Any]):
    issue_key = (payload.get("issueKey") or "").strip().upper()
    if not issue_key:
        return JSONResponse(status_code=400, content={"error": "issueKey is required"})
    url = f"{JIRA_BASE_URL}/rest/api/3/issue/{issue_key}"
    params = {"fields": "summary,description,status,assignee,issuetype,labels,priority,components,created,updated,reporter,parent"}
    response = requests.get(url, headers=jira_auth_headers(), params=params, timeout=30)
    if response.status_code != 200:
        return JSONResponse(status_code=response.status_code, content=safe_json_response(response))
    issue = response.json()
    normalized = normalize_live_issue(issue)
    cleaned = auto_clean_and_normalize_ticket(normalized["summary"], normalized["description"], normalized["issue_type"] or "Story")
    return {"issueKey": issue_key, "normalized_issue": normalized, "cleaned": cleaned}


@app.post("/api/tickets/bulk-dashboard")
def bulk_dashboard(payload: Dict[str, Any]):
    project_key = normalize_project_key(payload.get("projectKey"))
    exclude_done = bool(payload.get("excludeDoneCompleted", True))
    jql = payload.get("jql") or (f"project = {project_key} AND statusCategory != Done ORDER BY created DESC" if exclude_done else f"project = {project_key} ORDER BY created DESC")
    max_results = int(payload.get("maxResults") or 0)
    try:
        return build_bulk_dashboard(project_key, jql, max_results, exclude_done=exclude_done)
    except RuntimeError as e:
        return JSONResponse(status_code=500, content={"error": str(e)})



@app.post("/api/tickets/manager-dashboard")
def manager_dashboard(payload: Dict[str, Any]):
    project_key = normalize_project_key(payload.get("projectKey") or payload.get("project_key"))
    exclude_done = bool(payload.get("excludeDoneCompleted", True))
    jql = payload.get("jql") or (f"project = {project_key} AND statusCategory != Done ORDER BY created DESC" if exclude_done else f"project = {project_key} ORDER BY created DESC")
    max_results = int(payload.get("maxResults") or payload.get("max_results") or 0)
    try:
        return build_manager_dashboard(project_key, jql, max_results, exclude_done=exclude_done)
    except RuntimeError as e:
        return JSONResponse(status_code=500, content={"error": str(e)})

@app.post("/api/tickets/assignee-capacity")
def assignee_capacity(payload: Dict[str, Any]):
    assignee = (payload.get("assignee") or payload.get("assigneeName") or "").strip()
    assignee_account_id = (payload.get("assigneeAccountId") or payload.get("accountId") or "").strip()
    project_keys = parse_project_keys(payload.get("projectKeys") or payload.get("projects") or "")
    story_points_field_id = (payload.get("storyPointsFieldId") or payload.get("story_points_field_id") or "").strip()
    jql_override = (payload.get("jql") or payload.get("jqlOverride") or "").strip()
    auto_capacity = bool(payload.get("capacityAutoCalculated", False))
    raw_capacity_value = payload.get("sprintCapacityPoints") if payload.get("sprintCapacityPoints") not in (None, "") else payload.get("capacityPoints")
    if auto_capacity:
        sprint_capacity_points = 0.0
    else:
        sprint_capacity_points = float(raw_capacity_value or 20)
    max_results = int(payload.get("maxResults") or payload.get("max_results") or 0)
    start_date = (payload.get("startDate") or payload.get("capacityStartDate") or "").strip()
    end_date = (payload.get("endDate") or payload.get("capacityEndDate") or "").strip()
    if not assignee and not assignee_account_id:
        return JSONResponse(status_code=400, content={"error": "assignee or assigneeAccountId is required"})
    try:
        return build_capacity_dashboard(assignee, project_keys, story_points_field_id, sprint_capacity_points, jql_override, max_results, assignee_account_id, start_date, end_date)
    except ValueError as e:
        return JSONResponse(status_code=400, content={"error": str(e)})
    except RuntimeError as e:
        return JSONResponse(status_code=500, content={"error": str(e)})

@app.post("/api/tickets/team-capacity")
def team_capacity(payload: Dict[str, Any]):
    members = payload.get("members") or []
    project_keys = parse_project_keys(payload.get("projectKeys") or payload.get("projects") or "")
    story_points_field_id = (payload.get("storyPointsFieldId") or payload.get("story_points_field_id") or "").strip()
    jql_override = (payload.get("jql") or payload.get("jqlOverride") or "").strip()
    auto_capacity = bool(payload.get("capacityAutoCalculated", False))
    raw_default_capacity = payload.get("defaultCapacityPoints") if payload.get("defaultCapacityPoints") not in (None, "") else payload.get("capacityPoints")
    if auto_capacity:
        default_capacity_points = 0.0
    else:
        default_capacity_points = float(raw_default_capacity or 20)
    max_results = int(payload.get("maxResults") or payload.get("max_results") or 0)
    start_date = (payload.get("startDate") or payload.get("capacityStartDate") or "").strip()
    end_date = (payload.get("endDate") or payload.get("capacityEndDate") or "").strip()
    try:
        return build_team_capacity_dashboard(members, project_keys, story_points_field_id, default_capacity_points, jql_override, max_results, start_date, end_date)
    except ValueError as e:
        return JSONResponse(status_code=400, content={"error": str(e)})
    except RuntimeError as e:
        return JSONResponse(status_code=500, content={"error": str(e)})

@app.get("/api/tickets/bulk-dashboard.csv")
def bulk_dashboard_csv(projectKey: str = Query(default=DEFAULT_PROJECT), jql: Optional[str] = Query(default=None), maxResults: int = Query(default=0), excludeDoneCompleted: bool = Query(default=True)):
    project_key = normalize_project_key(projectKey)
    resolved_jql = jql or (f"project = {project_key} AND statusCategory != Done ORDER BY created DESC" if excludeDoneCompleted else f"project = {project_key} ORDER BY created DESC")
    result = build_bulk_dashboard(project_key, resolved_jql, maxResults, exclude_done=excludeDoneCompleted)
    output = csv_stream(bulk_results_to_csv_rows(result, weak_only=False))
    filename = f"{project_key.lower()}_ticket_quality_dashboard.csv"
    return StreamingResponse(iter([output.getvalue()]), media_type="text/csv", headers={"Content-Disposition": f'attachment; filename="{filename}"'})


@app.get("/api/tickets/bulk-dashboard-weak.csv")
def bulk_dashboard_weak_csv(projectKey: str = Query(default=DEFAULT_PROJECT), jql: Optional[str] = Query(default=None), maxResults: int = Query(default=0), excludeDoneCompleted: bool = Query(default=True)):
    project_key = normalize_project_key(projectKey)
    resolved_jql = jql or (f"project = {project_key} AND statusCategory != Done ORDER BY created DESC" if excludeDoneCompleted else f"project = {project_key} ORDER BY created DESC")
    result = build_bulk_dashboard(project_key, resolved_jql, maxResults, exclude_done=excludeDoneCompleted)
    output = csv_stream(bulk_results_to_csv_rows(result, weak_only=True))
    filename = f"{project_key.lower()}_weak_tickets.csv"
    return StreamingResponse(iter([output.getvalue()]), media_type="text/csv", headers={"Content-Disposition": f'attachment; filename="{filename}"'})


@app.post("/api/tickets/suggest-breakdown-live")
def suggest_breakdown_live(payload: Dict[str, Any]):
    issue_key = (payload.get("issueKey") or "").strip().upper()
    if not issue_key:
        return JSONResponse(status_code=400, content={"error": "issueKey is required"})
    url = f"{JIRA_BASE_URL}/rest/api/3/issue/{issue_key}"
    params = {"fields": "summary,description,issuetype,parent"}
    response = requests.get(url, headers=jira_auth_headers(), params=params, timeout=30)
    if response.status_code != 200:
        return JSONResponse(status_code=response.status_code, content=safe_json_response(response))
    issue = response.json()
    fields = issue.get("fields", {})
    summary = fields.get("summary") or ""
    description = extract_plain_text_from_adf(fields.get("description"))
    issue_type = (fields.get("issuetype") or {}).get("name", "")
    return suggest_breakdown(summary, description, issue_type, issue_key)


@app.post("/api/jira/create-batch")
def jira_create_batch(data: Dict[str, Any]):
    project_key = normalize_project_key(data.get("projectKey"))
    parent_issue_key = (data.get("parentIssueKey") or "").strip().upper()
    parent_issue_type = (data.get("parentIssueType") or "").strip()
    items = data.get("items") or []
    skip_existing = bool(data.get("skipExisting", True))
    story_points_field = (data.get("storyPointsFieldId") or JIRA_STORY_POINTS_FIELD or "").strip()

    created = []
    skipped = []
    errors = []

    for idx, item in enumerate(items):
        summary = (item.get("summary") or "").strip()
        description = (item.get("description") or "").strip()
        issue_type = (item.get("issueType") or "Task").strip()

        if not summary:
            errors.append({"index": idx, "error": "summary is required"})
            continue

        duplicate_check = find_existing_breakdown_ticket(project_key, parent_issue_key, issue_type, summary)
        if skip_existing and duplicate_check.get("found"):
            skipped.append({"index": idx, "action": "skipped", "reason": "Possible duplicate found under the same parent/project/type", "summary": summary, "issueType": issue_type, "existingIssueKey": duplicate_check.get("existingIssueKey"), "existingSummary": duplicate_check.get("existingSummary"), "duplicateDetectionMethod": duplicate_check.get("method")})
            continue

        app_labels = breakdown_labels(parent_issue_key, issue_type, summary)
        labels = dedupe_keep_order([*(item.get("labels") or []), *app_labels], 20)
        fields = {"project": {"key": project_key}, "summary": summary, "description": build_adf_document(description), "issuetype": {"name": issue_type}, "priority": {"name": item.get("priority") or "Medium"}, "labels": labels}
        components = item.get("components") or []
        if components:
            fields["components"] = [{"name": c} for c in components]
        if parent_issue_key:
            fields.update(parent_field_payload(parent_issue_key, parent_issue_type, issue_type))

        response = requests.post(f"{JIRA_BASE_URL}/rest/api/3/issue", headers=jira_auth_headers(), json={"fields": fields}, timeout=30)
        if response.status_code in (200, 201):
            create_body = safe_json_response(response)["jira_response"]
            created_key = create_body.get("key")
            story_points_result = None
            if issue_type.lower() == "story" and item.get("storyPoints") is not None and created_key:
                story_points_result = update_story_points_after_create(created_key, item.get("storyPoints"), story_points_field)
            created.append({"index": idx, "action": "created", "issueKey": created_key, "jira_response": create_body, "summary": summary, "issueType": issue_type, "labelsApplied": labels, "fingerprintLabel": f"fp-{breakdown_fingerprint(parent_issue_key, issue_type, summary)}", "storyPointsRequested": item.get("storyPoints") if issue_type.lower() == "story" else None, "storyPointsResult": story_points_result})
        else:
            errors.append({"index": idx, "status_code": response.status_code, "summary": summary, "issueType": issue_type, "jira_response": safe_json_response(response)["jira_response"], "duplicateCheck": duplicate_check})

    return {"projectKey": project_key, "parentIssueKey": parent_issue_key, "parentIssueType": parent_issue_type, "skipExisting": skip_existing, "storyPointsFieldPreferred": story_points_field or None, "created": created, "skipped": skipped, "errors": errors, "created_count": len(created), "skipped_count": len(skipped), "error_count": len(errors)}


@app.post("/api/jira/check-breakdown-duplicates")
def check_breakdown_duplicates(data: Dict[str, Any]):
    project_key = normalize_project_key(data.get("projectKey"))
    parent_issue_key = (data.get("parentIssueKey") or "").strip().upper()
    parent_issue_type = (data.get("parentIssueType") or "").strip()
    items = data.get("items") or []
    results = []
    for idx, item in enumerate(items):
        summary = (item.get("summary") or "").strip()
        issue_type = (item.get("issueType") or "Task").strip()
        if not summary:
            results.append({"index": idx, "found": False, "error": "summary is required"})
            continue
        result = find_existing_breakdown_ticket(project_key, parent_issue_key, issue_type, summary)
        result.update({"index": idx, "summary": summary, "issueType": issue_type})
        results.append(result)
    return {"projectKey": project_key, "parentIssueKey": parent_issue_key, "parentIssueType": parent_issue_type, "results": results}



@app.post("/api/tickets/enhance-with-docs")
async def enhance_with_docs(issueKey: str = Form(...), files: List[UploadFile] = File(...)):
    issue_key = (issueKey or "").strip().upper()
    if not issue_key:
        return JSONResponse(status_code=400, content={"error": "issueKey is required"})
    if not files:
        return JSONResponse(status_code=400, content={"error": "Upload at least one supporting document"})

    normalized, error = get_live_ticket_or_error(issue_key)
    if error:
        return error

    extracted_files = []
    for uploaded in files:
        filename = uploaded.filename or "uploaded_file"
        ext = get_file_extension(filename)
        content = await uploaded.read()
        if ext not in SUPPORTED_DOC_EXTENSIONS:
            extracted_files.append({"filename": filename, "extension": ext, "characters": 0, "text": f"Unsupported file type {ext}"})
            continue
        extracted_files.append(extract_text_from_uploaded_file(filename, content))

    doc_text = compact_doc_text(extracted_files)
    assessment = evaluate_ticket(normalized)
    deterministic_sections = infer_doc_sections(normalized, doc_text)
    ai_sections = call_openai_for_doc_enhancement(normalized, assessment, doc_text, deterministic_sections)
    final_sections = ai_sections or deterministic_sections
    for key, value in deterministic_sections.items():
        final_sections.setdefault(key, value)

    suggested_description = build_description_from_doc_sections(normalized.get("summary") or issue_key, final_sections)
    cleaned = auto_clean_and_normalize_ticket(normalized.get("summary") or issue_key, suggested_description, normalized.get("issue_type") or "Story")

    return {
        "issueKey": issue_key,
        "normalized_issue": normalized,
        "assessment_before": assessment,
        "files_processed": [{"filename": f.get("filename"), "extension": f.get("extension"), "characters": f.get("characters")} for f in extracted_files],
        "used_openai": bool(ai_sections),
        "mandatory_criteria_suggestions": final_sections,
        "suggested_description": cleaned.get("normalized_description") or suggested_description,
        "suggested_summary": cleaned.get("summary") or normalized.get("summary"),
    }

@app.post("/api/tickets/apply-doc-suggestions")
def apply_doc_suggestions(data: Dict[str, Any]):
    issue_key = (data.get("issueKey") or "").strip().upper()
    description = (data.get("description") or data.get("suggested_description") or "").strip()
    summary = data.get("summary") or data.get("suggested_summary")
    if not issue_key:
        return JSONResponse(status_code=400, content={"error": "issueKey is required"})
    if not description:
        return JSONResponse(status_code=400, content={"error": "description/suggested_description is required"})
    payload = {"issueKey": issue_key, "description": description}
    if summary and str(summary).strip():
        payload["summary"] = str(summary).strip()
    result = jira_update(payload)
    if isinstance(result, JSONResponse):
        return result
    result["appliedDocSuggestions"] = True
    return result


@app.post("/api/jira/create")
def jira_create(data: Dict[str, Any]):
    project_key = normalize_project_key(data.get("projectKey"))
    issue_type = data.get("issueType", "Story")
    summary = data.get("summary", "").strip()
    description = data.get("description", "").strip()
    if not summary:
        return JSONResponse(status_code=400, content={"error": "summary is required"})
    payload = {"fields": {"project": {"key": project_key}, "summary": summary, "description": build_adf_document(description), "issuetype": {"name": issue_type}}}
    response = requests.post(f"{JIRA_BASE_URL}/rest/api/3/issue", headers=jira_auth_headers(), json=payload, timeout=30)
    result = safe_json_response(response)
    result["projectKey"] = project_key
    return result


@app.post("/api/jira/update")
def jira_update(data: Dict[str, Any]):
    issue_key = (data.get("issueKey") or "").strip().upper()
    summary = data.get("summary")
    description = data.get("description")
    issue_type = data.get("issueType")
    if not issue_key:
        return JSONResponse(status_code=400, content={"error": "issueKey is required"})
    fields_to_update: Dict[str, Any] = {}
    if summary is not None and str(summary).strip():
        fields_to_update["summary"] = str(summary).strip()
    if description is not None:
        fields_to_update["description"] = build_adf_document(str(description))
    if issue_type is not None and str(issue_type).strip():
        fields_to_update["issuetype"] = {"name": str(issue_type).strip()}
    if not fields_to_update:
        return JSONResponse(status_code=400, content={"error": "At least one of summary, description, or issueType must be provided"})
    response = requests.put(f"{JIRA_BASE_URL}/rest/api/3/issue/{issue_key}", headers=jira_auth_headers(), json={"fields": fields_to_update}, timeout=30)
    result = {"status_code": response.status_code, "issueKey": issue_key}
    if response.text.strip():
        try:
            result["jira_response"] = response.json()
        except Exception:
            result["jira_response"] = {"raw_text": response.text}
    else:
        result["jira_response"] = {"message": "Issue updated successfully"}
    return result
@app.get("/health")
def health():
    return {"status": "ok"}
